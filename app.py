"""
Convin Klaro
AI-powered knowledge & support intelligence platform
"""

import streamlit as st
import json, re, os, base64, zipfile, io
from datetime import datetime
from urllib.parse import urljoin, urlparse
import anthropic

def _b64_img(path):
    """Return base64 data-URI for an image file, or '' if not found."""
    try:
        with open(path, "rb") as f:
            ext = path.rsplit(".", 1)[-1].lower()
            mime = {"png":"image/png","jpg":"image/jpeg","jpeg":"image/jpeg",
                    "svg":"image/svg+xml","webp":"image/webp"}.get(ext,"image/png")
            return f"data:{mime};base64,{base64.b64encode(f.read()).decode()}"
    except Exception:
        return ""

_LOGO_URI = _b64_img(os.path.join(os.path.dirname(os.path.abspath(__file__)), "convin_logo.png"))

# ══════════════════════════════════════════════════════════════════
#  CONFIG
# ══════════════════════════════════════════════════════════════════
APP_DIR = os.path.dirname(os.path.abspath(__file__))
KB_FILE = os.path.join(APP_DIR, "kb_store.json")
KB_KEYS = ("kb_documents", "kb_links", "kb_whatsapp", "kb_crawled", "kb_faqs")
MAX_CTX  = 580_000   # chars — safely under Claude's 200k-token window

st.set_page_config(
    page_title="Convin Klaro",
    page_icon="✦",
    layout="wide",
    initial_sidebar_state="collapsed",
)

# ══════════════════════════════════════════════════════════════════
#  GLOBAL STYLES
# ══════════════════════════════════════════════════════════════════
st.markdown("""
<style>
@import url('https://fonts.googleapis.com/css2?family=Inter:wght@300;400;500;600;700;800&display=swap');

/* ── Hide Streamlit chrome ── */
#MainMenu, footer, header,
[data-testid="stSidebar"],
[data-testid="collapsedControl"]   { display: none !important; }
.stDeployButton                    { display: none !important; }
[data-testid="stToolbar"]          { display: none !important; }

/*
 * ═══════════════════════════════════════════════════════════
 *  CONVIN KLARO — Premium Calm Design System
 * ───────────────────────────────────────────────────────────
 *  BG:         #0B0F1A   deep navy base
 *  Surface-1:  #111827   glass cards
 *  Surface-2:  #1A2035   elevated panels
 *  Surface-3:  #0D1117   deep answer blocks
 *  Accent:     #6366F1   indigo (calm, trustworthy)
 *  Accent-lt:  #818CF8   light indigo
 *  Cyan:       #22D3EE   subtle highlight
 *  Text-1:     #E5E7EB   soft white
 *  Text-2:     #94A3B8   muted steel blue
 *  Text-3:     #475569   very muted
 *  Border:     rgba(255,255,255,0.06)
 *  Border-acc: rgba(99,102,241,0.22)
 *  Green:      #10B981   WhatsApp / live
 * ═══════════════════════════════════════════════════════════
 */

/* ── Base ── */
html, body, [class*="css"] {
    font-family: 'Inter', -apple-system, BlinkMacSystemFont, sans-serif;
    -webkit-font-smoothing: antialiased;
    -moz-osx-font-smoothing: grayscale;
}
.main {
    background: #060914;
    background-image:
        radial-gradient(ellipse 110% 55% at 50% -10%, rgba(139,92,246,0.18) 0%, transparent 55%),
        radial-gradient(ellipse 70%  45% at 95%  90%, rgba(236,72,153,0.12) 0%, transparent 50%),
        radial-gradient(ellipse 60%  40% at 5%   75%, rgba(34,211,238,0.09) 0%, transparent 50%),
        radial-gradient(ellipse 55%  35% at 75%  20%, rgba(99,102,241,0.10) 0%, transparent 52%);
}
.main .block-container { padding: 0 !important; max-width: 100% !important; }

/* Aurora mesh grid */
[data-testid="stAppViewContainer"]::before {
    content: '';
    position: fixed; top: 0; left: 0; right: 0; bottom: 0;
    background-image:
        radial-gradient(rgba(139,92,246,0.055) 1px, transparent 1px),
        radial-gradient(rgba(34,211,238,0.025) 1px, transparent 1px);
    background-size: 38px 38px, 72px 72px;
    background-position: 0 0, 19px 19px;
    pointer-events: none; z-index: 0;
    animation: gridshift 50s linear infinite;
}
@keyframes gridshift {
    0%   { background-position: 0 0, 19px 19px; }
    100% { background-position: 38px 38px, 57px 57px; }
}

/* Multi-colour aurora orbs */
[data-testid="stAppViewContainer"]::after {
    content: '';
    position: fixed; top: -30%; left: -20%;
    width: 80%; height: 80%;
    background:
        radial-gradient(ellipse at 30% 30%, rgba(139,92,246,0.12) 0%, transparent 55%),
        radial-gradient(ellipse at 70% 70%, rgba(236,72,153,0.08) 0%, transparent 55%),
        radial-gradient(ellipse at 80% 10%, rgba(34,211,238,0.07) 0%, transparent 50%);
    border-radius: 50%;
    pointer-events: none; z-index: 0;
    animation: aurora 28s ease-in-out infinite alternate;
}
@keyframes aurora {
    0%   { transform: translate(0%,   0%)   scale(1)    rotate(0deg);   }
    33%  { transform: translate(22%,  16%)  scale(1.10) rotate(4deg);   }
    66%  { transform: translate(8%,   32%)  scale(0.94) rotate(-3deg);  }
    100% { transform: translate(28%,  10%)  scale(1.06) rotate(2deg);   }
}

/* ══ TOP NAV ══════════════════════════════════════════════ */
.topnav {
    position: fixed; top: 0; left: 0; right: 0; z-index: 999;
    height: 62px;
    background: rgba(6,9,20,0.88);
    backdrop-filter: blur(32px) saturate(220%);
    -webkit-backdrop-filter: blur(32px) saturate(220%);
    border-bottom: 1px solid rgba(139,92,246,0.18);
    display: flex; align-items: center;
    padding: 0 40px;
    justify-content: space-between;
    box-shadow: 0 1px 0 rgba(139,92,246,0.12), 0 4px 32px rgba(0,0,0,0.30);
}
.topnav-brand { display: flex; align-items: center; gap: 12px; }
.topnav-logo  { height: 28px; width: auto; object-fit: contain; opacity: 0.92; }
.topnav-brand .dot {
    width: 34px; height: 34px; border-radius: 10px;
    background: linear-gradient(135deg, #8B5CF6 0%, #6366F1 50%, #EC4899 100%);
    display: flex; align-items: center; justify-content: center;
    font-size: 0.88rem; color: #fff; font-weight: 700;
    box-shadow: 0 0 22px rgba(139,92,246,0.45), 0 4px 12px rgba(0,0,0,0.3);
}
.topnav-brand .name {
    font-size: 0.95rem; font-weight: 700; color: #E5E7EB;
    letter-spacing: -0.025em;
}
.topnav-brand .badge {
    font-size: 0.57rem; font-weight: 700; letter-spacing: 0.11em;
    background: linear-gradient(135deg, rgba(139,92,246,0.15), rgba(236,72,153,0.10));
    color: #C4B5FD;
    border: 1px solid rgba(139,92,246,0.32);
    padding: 3px 12px; border-radius: 20px;
    text-transform: uppercase;
    box-shadow: 0 0 12px rgba(139,92,246,0.12);
}
.topnav-right { display: flex; align-items: center; gap: 10px; }
.topnav-status {
    display: flex; align-items: center; gap: 6px;
    font-size: 0.71rem; color: #9CA3AF; font-weight: 500;
}
.live-dot {
    width: 7px; height: 7px; border-radius: 50%;
    background: #10B981;
    box-shadow: 0 0 8px rgba(16,185,129,0.60);
    animation: livepulse 3s ease-in-out infinite;
}
@keyframes livepulse {
    0%,100% { opacity:1; transform:scale(1); }
    50%      { opacity:0.55; transform:scale(0.78); }
}

/* ══ CHAT PAGE ════════════════════════════════════════════ */
.chat-page { min-height:100vh; padding-top:62px; }
.chat-feed  { max-width:740px; width:100%; margin:0 auto; padding:48px 24px 180px; }

/* Welcome card */
.welcome-card { text-align:center; padding:82px 24px 54px; }
.welcome-icon {
    width: 72px; height: 72px; border-radius: 24px;
    background: linear-gradient(135deg, #6366F1 0%, #4338CA 100%);
    margin: 0 auto 26px;
    display: flex; align-items: center; justify-content: center;
    font-size: 2rem;
    box-shadow:
        0 0 52px rgba(99,102,241,0.28),
        0 14px 36px rgba(0,0,0,0.42),
        inset 0 1px 0 rgba(255,255,255,0.14);
    animation: float 5.5s ease-in-out infinite;
}
@keyframes float {
    0%,100% { transform: translateY(0) rotate(0deg);  }
    50%      { transform: translateY(-9px) rotate(1deg); }
}
.welcome-eyebrow {
    display: inline-block;
    font-size: 0.64rem; font-weight: 600; letter-spacing: 0.14em;
    text-transform: uppercase; color: #818CF8;
    background: rgba(99,102,241,0.10);
    border: 1px solid rgba(99,102,241,0.22);
    padding: 4px 14px; border-radius: 20px;
    margin-bottom: 18px;
}
.welcome-title {
    font-size: 1.78rem; font-weight: 800;
    color: #E5E7EB; letter-spacing: -0.03em;
    line-height: 1.18; margin-bottom: 13px;
}
.welcome-title span {
    background: linear-gradient(90deg, #818CF8 0%, #22D3EE 100%);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
}
.welcome-sub {
    font-size: 0.88rem; color: #94A3B8; line-height: 1.72;
    max-width: 420px; margin: 0 auto;
}

/* KB stat chips */
.kb-stats-row { display:flex; gap:10px; justify-content:center; flex-wrap:wrap; margin-top:28px; }
.kb-stat-chip {
    display: flex; align-items: center; gap: 8px;
    background: rgba(17,24,39,0.65);
    backdrop-filter: blur(14px);
    border: 1px solid rgba(255,255,255,0.06);
    border-radius: 30px; padding: 8px 18px;
    font-size: 0.78rem; color: #9CA3AF;
    animation: popIn 0.42s ease forwards; opacity: 0;
    transition: border-color 0.2s, background 0.2s;
}
.kb-stat-chip:nth-child(1) { animation-delay: 0.10s; }
.kb-stat-chip:nth-child(2) { animation-delay: 0.22s; }
.kb-stat-chip:nth-child(3) { animation-delay: 0.34s; }
@keyframes popIn {
    from { opacity:0; transform:translateY(8px) scale(0.95); }
    to   { opacity:1; transform:translateY(0)    scale(1);    }
}
.kb-stat-chip.active {
    border-color: rgba(99,102,241,0.26);
    background: rgba(99,102,241,0.08);
    color: #6B7280;
}
.kb-stat-chip.wa.active {
    border-color: rgba(16,185,129,0.26);
    background: rgba(16,185,129,0.07);
}
.kb-stat-chip .num {
    font-weight: 800; font-size: 0.9rem;
    background: linear-gradient(135deg, #818CF8, #6366F1);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
}
.kb-stat-chip.wa .num {
    background: linear-gradient(135deg, #34D399, #10B981);
    -webkit-background-clip: text; -webkit-text-fill-color: transparent;
}
.no-kb-hint {
    margin-top: 22px; font-size: 0.76rem; color: #6B7280;
    background: rgba(255,255,255,0.02);
    border: 1px dashed rgba(255,255,255,0.07);
    border-radius: 20px; padding: 8px 22px; display: inline-block;
}
.ready-badge {
    display: inline-flex; align-items: center; gap: 7px;
    margin-top: 18px;
    font-size: 0.63rem; font-weight: 700; letter-spacing: 0.13em;
    text-transform: uppercase; color: #10B981;
    background: rgba(16,185,129,0.08);
    border: 1px solid rgba(16,185,129,0.22);
    padding: 4px 16px; border-radius: 20px;
    animation: readypop 0.6s 0.5s ease forwards; opacity: 0;
}
@keyframes readypop {
    from { opacity:0; transform:scale(0.92); }
    to   { opacity:1; transform:scale(1);    }
}

/* Suggestion chips */
[data-testid="stHorizontalBlock"] .stButton > button[kind="secondary"] {
    border-radius: 24px !important;
    padding: 9px 16px !important;
    font-size: 0.79rem !important;
    background: rgba(17,24,39,0.65) !important;
    backdrop-filter: blur(14px) !important;
    border: 1px solid rgba(255,255,255,0.10) !important;
    color: #9CA3AF !important;
    white-space: nowrap !important;
    transition: all 0.25s cubic-bezier(0.4,0,0.2,1) !important;
}
[data-testid="stHorizontalBlock"] .stButton > button[kind="secondary"]:hover {
    border-color: rgba(99,102,241,0.38) !important;
    color: #818CF8 !important;
    background: rgba(99,102,241,0.08) !important;
    transform: translateY(-2px) !important;
    box-shadow: 0 8px 26px rgba(99,102,241,0.14) !important;
}

/* Messages */
.msg-group { margin-bottom: 28px; animation: fadeUp 0.26s cubic-bezier(0.4,0,0.2,1); }
@keyframes fadeUp {
    from { opacity:0; transform:translateY(12px); }
    to   { opacity:1; transform:translateY(0);    }
}
.msg-user-row { display:flex; justify-content:flex-end; margin-bottom:5px; }
.msg-user-bubble {
    background: linear-gradient(135deg, #6366F1 0%, #4338CA 100%);
    color: #EEF2FF;
    padding: 13px 20px;
    border-radius: 20px 20px 5px 20px;
    max-width: 65%;
    font-size: 0.88rem; line-height: 1.65;
    box-shadow: 0 4px 26px rgba(99,102,241,0.28), 0 2px 8px rgba(0,0,0,0.22);
}
.msg-ai-row { display:flex; align-items:flex-start; gap:12px; }
.msg-ai-avatar {
    position: relative;
    width: 34px; height: 34px; border-radius: 11px;
    background: linear-gradient(135deg, #6366F1 0%, #4338CA 100%);
    display: flex; align-items: center; justify-content: center;
    font-size: 0.72rem; color: #EEF2FF; font-weight: 700;
    flex-shrink: 0; margin-top: 2px;
    box-shadow: 0 0 18px rgba(99,102,241,0.28), 0 4px 10px rgba(0,0,0,0.25);
}
.msg-ai-avatar::before {
    content: '';
    position: absolute; top: -2px; left: -2px;
    width: 38px; height: 38px; border-radius: 13px;
    border: 1.5px solid rgba(99,102,241,0.25);
    animation: avatarring 4s ease-out infinite;
}
@keyframes avatarring {
    0%   { opacity:0.8; transform:scale(1);    }
    100% { opacity:0;   transform:scale(1.55); }
}
.msg-ai-bubble {
    background: rgba(17,24,39,0.72);
    backdrop-filter: blur(18px);
    -webkit-backdrop-filter: blur(18px);
    color: #94A3B8;
    padding: 14px 20px;
    border-radius: 5px 20px 20px 20px;
    max-width: 72%;
    font-size: 0.88rem; line-height: 1.80;
    border: 1px solid rgba(255,255,255,0.06);
    box-shadow: 0 4px 22px rgba(0,0,0,0.20);
}
.msg-ai-bubble b, .msg-ai-bubble strong { color: #818CF8; font-weight: 600; }
.msg-ai-bubble ol, .msg-ai-bubble ul    { padding-left: 18px; margin: 8px 0; }
.msg-ai-bubble li { margin-bottom: 5px; color: #94A3B8; }
.msg-ai-bubble blockquote {
    border-left: 3px solid #10B981;
    background: rgba(16,185,129,0.07);
    border-radius: 0 9px 9px 0;
    padding: 7px 14px; margin: 10px 0 4px;
    font-size: 0.82rem; color: #34D399;
}
.msg-ts       { font-size:0.6rem; color:#6B7280; margin-top:5px; text-align:right; }
.msg-ts-left  { text-align:left; margin-left:46px; }
.sources-bar  { display:flex; align-items:center; gap:5px; margin-top:8px; margin-left:46px; flex-wrap:wrap; }
.source-tag {
    font-size: 0.62rem; font-weight: 500; color: #6B7280;
    background: rgba(255,255,255,0.04);
    border: 1px solid rgba(255,255,255,0.08);
    padding: 2px 9px; border-radius: 10px;
}

/* Input bar */
.input-bar-wrap {
    position: fixed; bottom: 0; left: 0; right: 0; z-index: 998;
    background: rgba(11,15,26,0.88);
    backdrop-filter: blur(28px) saturate(180%);
    -webkit-backdrop-filter: blur(28px) saturate(180%);
    border-top: 1px solid rgba(255,255,255,0.05);
    padding: 14px 0 22px;
    box-shadow: 0 -1px 30px rgba(0,0,0,0.28);
}
.input-bar-inner { max-width:740px; margin:0 auto; padding:0 24px; display:flex; align-items:center; gap:10px; }
.input-bar-inner .stTextInput > div > div {
    background: rgba(17,24,39,0.78) !important;
    backdrop-filter: blur(16px) !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    border-radius: 16px !important;
    font-size: 0.9rem !important; color: #E5E7EB !important;
    transition: border-color 0.2s, box-shadow 0.2s !important;
    box-shadow: 0 2px 12px rgba(0,0,0,0.16) !important;
}
.input-bar-inner .stTextInput > div > div:focus-within {
    border-color: rgba(99,102,241,0.52) !important;
    box-shadow: 0 0 0 3px rgba(99,102,241,0.09), 0 2px 12px rgba(0,0,0,0.16) !important;
}
.input-bar-inner .stTextInput input { font-size:0.9rem !important; color:#E5E7EB !important; background:transparent !important; }
.input-bar-inner .stTextInput input::placeholder { color:#4B5563 !important; }

/* ══ BUTTONS ════════════════════════════════════════════ */
.stButton > button {
    border-radius: 11px !important; font-weight: 600 !important;
    font-size: 0.83rem !important;
    transition: all 0.22s cubic-bezier(0.4,0,0.2,1) !important;
    letter-spacing: -0.01em !important;
}
.stButton > button[kind="primary"] {
    background: linear-gradient(135deg, #8B5CF6 0%, #6366F1 50%, #4F46E5 100%) !important;
    color: #EEF2FF !important; border: none !important;
    box-shadow: 0 4px 24px rgba(139,92,246,0.42), 0 2px 8px rgba(0,0,0,0.25) !important;
}
.stButton > button[kind="primary"]:hover {
    background: linear-gradient(135deg, #A78BFA 0%, #818CF8 50%, #6366F1 100%) !important;
    box-shadow: 0 8px 36px rgba(139,92,246,0.55), 0 2px 10px rgba(0,0,0,0.25) !important;
    transform: translateY(-2px) !important;
}
.stButton > button[kind="secondary"] {
    background: rgba(17,24,39,0.65) !important;
    color: #9CA3AF !important;
    border: 1px solid rgba(255,255,255,0.10) !important;
    backdrop-filter: blur(10px) !important;
}
.stButton > button[kind="secondary"]:hover {
    border-color: rgba(99,102,241,0.38) !important;
    color: #C4C9D4 !important;
    background: rgba(99,102,241,0.08) !important;
    box-shadow: 0 4px 16px rgba(99,102,241,0.12) !important;
}

/* ══ TEXT INPUTS ════════════════════════════════════════ */
.stTextInput > div > div {
    background: rgba(17,24,39,0.68) !important;
    border: 1px solid rgba(255,255,255,0.08) !important;
    border-radius: 12px !important; color: #E5E7EB !important;
    transition: border-color 0.2s, box-shadow 0.2s !important;
    backdrop-filter: blur(12px) !important;
}
.stTextInput > div > div:focus-within {
    border-color: rgba(99,102,241,0.48) !important;
    box-shadow: 0 0 0 3px rgba(99,102,241,0.09) !important;
}
.stTextInput input { color: #E5E7EB !important; background: transparent !important; }
.stTextInput input::placeholder { color: #4B5563 !important; }
.stTextInput label { color: #94A3B8 !important; font-size: 0.82rem !important; font-weight: 500 !important; }

/* ══ SETTINGS PAGE ═══════════════════════════════════════ */
.settings-page { padding-top: 62px; min-height: 100vh; }
.settings-inner { max-width: 900px; margin: 0 auto; padding: 44px 32px; }
.settings-title { font-size:1.3rem; font-weight:800; color:#E5E7EB; letter-spacing:-0.025em; margin-bottom:4px; }
.settings-sub   { font-size:0.82rem; color:#9CA3AF; margin-bottom:30px; }

/* File rows */
.file-row {
    display: flex; align-items: center; gap: 12px;
    padding: 11px 16px; border-radius: 12px;
    background: rgba(17,24,39,0.60);
    backdrop-filter: blur(12px);
    border: 1px solid rgba(255,255,255,0.06);
    margin-bottom: 7px;
    transition: border-color 0.2s, background 0.2s;
}
.file-row:hover {
    border-color: rgba(99,102,241,0.24);
    background: rgba(17,24,39,0.80);
}
.file-row-icon { font-size:1.05rem; flex-shrink:0; }
.file-row-name { font-size:0.82rem; font-weight:500; color:#9CA3AF; flex:1; }
.file-row-meta { font-size:0.69rem; color:#6B7280; }

/* Upload zone */
div[data-testid="stFileUploader"] {
    background: rgba(99,102,241,0.04) !important;
    border: 1.5px dashed rgba(99,102,241,0.25) !important;
    border-radius: 14px !important;
    transition: border-color 0.2s, background 0.2s !important;
}
div[data-testid="stFileUploader"]:hover {
    background: rgba(99,102,241,0.07) !important;
    border-color: rgba(99,102,241,0.42) !important;
}

/* ══ TABS ════════════════════════════════════════════════ */
.stTabs [data-baseweb="tab-list"] {
    background: rgba(10,8,25,0.70);
    backdrop-filter: blur(18px);
    border-radius: 14px; padding: 5px; gap: 3px;
    border: 1px solid rgba(139,92,246,0.18);
    box-shadow: 0 4px 24px rgba(0,0,0,0.25), inset 0 1px 0 rgba(255,255,255,0.04);
}
.stTabs [data-baseweb="tab"] {
    border-radius: 10px !important; font-size: 0.8rem !important;
    font-weight: 600 !important; padding: 8px 18px !important;
    color: #6B7280 !important;
    transition: color 0.2s, background 0.2s !important;
}
.stTabs [data-baseweb="tab"]:hover { color: #C4B5FD !important; background: rgba(139,92,246,0.07) !important; }
.stTabs [aria-selected="true"] {
    background: linear-gradient(135deg, rgba(139,92,246,0.22), rgba(99,102,241,0.16)) !important;
    color: #C4B5FD !important;
    box-shadow: 0 0 0 1px rgba(139,92,246,0.35), 0 4px 14px rgba(139,92,246,0.18) !important;
}

/* Toggle */
.stCheckbox label { font-size:0.85rem !important; color:#94A3B8 !important; font-weight:500 !important; }

/* Slider */
.stSlider [data-baseweb="slider"] div[role="slider"] {
    background: #6366F1 !important;
    box-shadow: 0 0 10px rgba(99,102,241,0.50) !important;
}

/* Divider */
hr { border: none; border-top: 1px solid rgba(255,255,255,0.05) !important; margin: 22px 0 !important; }

/* ══ EXPANDERS ═══════════════════════════════════════════ */
.streamlit-expanderHeader {
    background: rgba(17,24,39,0.68) !important;
    backdrop-filter: blur(14px) !important;
    border: 1px solid rgba(255,255,255,0.07) !important;
    border-radius: 12px !important;
    color: #6B7280 !important;
    font-size: 0.88rem !important; font-weight: 600 !important;
    padding: 12px 16px !important;
    transition: border-color 0.2s, color 0.2s !important;
}
.streamlit-expanderHeader:hover {
    border-color: rgba(99,102,241,0.28) !important;
    color: #94A3B8 !important;
}
.streamlit-expanderContent {
    background: rgba(11,15,26,0.80) !important;
    border: 1px solid rgba(255,255,255,0.05) !important;
    border-top: none !important;
    border-radius: 0 0 12px 12px !important;
}
[data-testid="stExpander"] {
    border: 1px solid rgba(255,255,255,0.07) !important;
    border-radius: 12px !important;
    background: rgba(17,24,39,0.62) !important;
    backdrop-filter: blur(14px) !important;
    -webkit-backdrop-filter: blur(14px) !important;
    margin-bottom: 9px !important;
    overflow: hidden !important;
    transition: border-color 0.22s, box-shadow 0.22s !important;
}
[data-testid="stExpander"]:hover {
    border-color: rgba(99,102,241,0.26) !important;
    box-shadow: 0 4px 28px rgba(0,0,0,0.16) !important;
}
[data-testid="stExpander"] summary {
    background: transparent !important;
    color: #9CA3AF !important;
    font-size: 0.88rem !important; font-weight: 600 !important;
    border-radius: 12px !important;
    padding: 13px 18px !important;
    transition: color 0.2s !important;
}
[data-testid="stExpander"] summary:hover { color: #E5E7EB !important; }
[data-testid="stExpander"] summary span,
[data-testid="stExpander"] summary p { color: inherit !important; font-weight: 600 !important; }
[data-testid="stExpander"] > div:last-child {
    background: rgba(11,15,26,0.75) !important;
    border-top: 1px solid rgba(255,255,255,0.05) !important;
    padding: 0 !important;
}
[data-testid="stExpander"] > div:last-child *,
[data-testid="stExpander"] > div:last-child p,
[data-testid="stExpander"] > div:last-child span,
[data-testid="stExpander"] > div:last-child li,
[data-testid="stExpander"] > div:last-child div { color: #94A3B8 !important; }
[data-testid="stExpander"] > div:last-child strong,
[data-testid="stExpander"] > div:last-child b    { color: #818CF8 !important; }

/* Progress */
.stProgress > div > div {
    background: linear-gradient(90deg, #6366F1, #22D3EE) !important;
    border-radius: 4px !important;
}

/* Metrics */
[data-testid="metric-container"] {
    background: rgba(17,24,39,0.62);
    backdrop-filter: blur(14px);
    border: 1px solid rgba(255,255,255,0.06);
    border-radius: 14px; padding: 16px 20px;
    box-shadow: 0 2px 14px rgba(0,0,0,0.15);
    transition: border-color 0.2s, box-shadow 0.2s;
}
[data-testid="metric-container"]:hover {
    border-color: rgba(99,102,241,0.20);
    box-shadow: 0 4px 24px rgba(0,0,0,0.20);
}
[data-testid="metric-container"] [data-testid="stMetricValue"] {
    color: #E5E7EB !important; font-weight: 800 !important; font-size: 1.6rem !important;
}
[data-testid="metric-container"] [data-testid="stMetricLabel"] {
    color: #9CA3AF !important; font-size: 0.75rem !important; font-weight: 500 !important;
}

/* Alerts */
.stAlert { border-radius: 12px !important; border: none !important; }

/* Scrollbar */
::-webkit-scrollbar { width: 5px; }
::-webkit-scrollbar-track { background: transparent; }
::-webkit-scrollbar-thumb { background: rgba(99,102,241,0.16); border-radius: 5px; }
::-webkit-scrollbar-thumb:hover { background: rgba(99,102,241,0.32); }

/* ══ LANDING PAGE HERO ════════════════════════════════════ */
.lp-hero {
    position: relative; overflow: hidden;
    padding: 52px 44px 44px;
    margin-bottom: 32px;
    border-radius: 28px;
    background:
        linear-gradient(135deg,
            rgba(17,7,40,0.95)  0%,
            rgba(26,14,58,0.95) 35%,
            rgba(14,22,50,0.95) 65%,
            rgba(6,12,30,0.95)  100%);
    border: 1px solid rgba(139,92,246,0.30);
    box-shadow:
        0 0  0  1px rgba(236,72,153,0.08),
        0 24px 80px rgba(0,0,0,0.55),
        inset 0 1px 0  rgba(255,255,255,0.07);
}

/* Animated gradient border shimmer */
.lp-hero::before {
    content: '';
    position: absolute; inset: -2px;
    border-radius: 30px;
    background: linear-gradient(130deg,
        rgba(139,92,246,0.55) 0%,
        rgba(236,72,153,0.40) 30%,
        rgba(34,211,238,0.35) 55%,
        rgba(245,158,11,0.30) 75%,
        rgba(139,92,246,0.55) 100%);
    background-size: 300% 300%;
    animation: bordershine 6s ease infinite;
    z-index: -1;
}
@keyframes bordershine {
    0%   { background-position: 0%   50%; }
    50%  { background-position: 100% 50%; }
    100% { background-position: 0%   50%; }
}

/* Floating glow orbs inside hero */
.lp-hero::after {
    content: '';
    position: absolute; top: -60px; right: -40px;
    width: 320px; height: 320px; border-radius: 50%;
    background: radial-gradient(circle,
        rgba(139,92,246,0.22) 0%,
        rgba(236,72,153,0.10) 45%,
        transparent 70%);
    pointer-events: none;
    animation: heroorb 8s ease-in-out infinite alternate;
}
@keyframes heroorb {
    0%   { transform: translate(0, 0)    scale(1);    }
    100% { transform: translate(-30px, 30px) scale(1.15); }
}

.lp-hero-inner { position: relative; z-index: 1; }

/* Eyebrow pill */
.lp-eyebrow {
    display: inline-flex; align-items: center; gap: 8px;
    font-size: 0.65rem; font-weight: 700; letter-spacing: 0.14em;
    text-transform: uppercase;
    background: linear-gradient(135deg, rgba(139,92,246,0.18), rgba(236,72,153,0.12));
    border: 1px solid rgba(139,92,246,0.40);
    color: #C4B5FD;
    padding: 5px 16px; border-radius: 30px;
    margin-bottom: 20px;
    box-shadow: 0 0 20px rgba(139,92,246,0.15);
}
.lp-eyebrow .dot {
    width: 6px; height: 6px; border-radius: 50%;
    background: linear-gradient(135deg, #A78BFA, #EC4899);
    box-shadow: 0 0 8px rgba(167,139,250,0.8);
    animation: livepulse 2.5s ease-in-out infinite;
}

/* Main title */
.lp-title {
    font-size: clamp(1.8rem, 3.5vw, 2.6rem) !important;
    font-weight: 900 !important;
    line-height: 1.12 !important;
    letter-spacing: -0.04em !important;
    color: #F1F5F9 !important;
    margin: 0 0 14px !important;
}
.lp-title .grad {
    background: linear-gradient(95deg,
        #A78BFA 0%,
        #EC4899 35%,
        #22D3EE 65%,
        #A78BFA 100%);
    background-size: 200% auto;
    -webkit-background-clip: text;
    -webkit-text-fill-color: transparent;
    background-clip: text;
    animation: textshine 4s linear infinite;
}
@keyframes textshine {
    0%   { background-position: 0%   center; }
    100% { background-position: 200% center; }
}

/* Subtitle */
.lp-sub {
    font-size: 0.92rem !important;
    color: #94A3B8 !important;
    line-height: 1.7 !important;
    max-width: 540px;
    margin-bottom: 28px !important;
}

/* Feature pills row */
.lp-pills {
    display: flex; flex-wrap: wrap; gap: 8px;
    margin-bottom: 32px;
}
.lp-pill {
    display: inline-flex; align-items: center; gap: 6px;
    font-size: 0.72rem; font-weight: 600;
    padding: 6px 16px; border-radius: 30px;
    letter-spacing: 0.02em;
    transition: transform 0.2s, box-shadow 0.2s;
    cursor: default;
}
.lp-pill:hover { transform: translateY(-2px); }
.lp-pill-violet {
    background: rgba(139,92,246,0.14);
    border: 1px solid rgba(139,92,246,0.40);
    color: #C4B5FD;
    box-shadow: 0 2px 12px rgba(139,92,246,0.15);
}
.lp-pill-pink {
    background: rgba(236,72,153,0.12);
    border: 1px solid rgba(236,72,153,0.38);
    color: #F9A8D4;
    box-shadow: 0 2px 12px rgba(236,72,153,0.13);
}
.lp-pill-cyan {
    background: rgba(34,211,238,0.10);
    border: 1px solid rgba(34,211,238,0.35);
    color: #67E8F9;
    box-shadow: 0 2px 12px rgba(34,211,238,0.12);
}
.lp-pill-green {
    background: rgba(16,185,129,0.10);
    border: 1px solid rgba(16,185,129,0.35);
    color: #6EE7B7;
    box-shadow: 0 2px 12px rgba(16,185,129,0.12);
}
.lp-pill-amber {
    background: rgba(245,158,11,0.10);
    border: 1px solid rgba(245,158,11,0.35);
    color: #FCD34D;
    box-shadow: 0 2px 12px rgba(245,158,11,0.12);
}

/* Stat cards row */
.lp-stats {
    display: flex; gap: 14px; flex-wrap: wrap;
}
.lp-stat {
    flex: 1; min-width: 100px;
    background: rgba(255,255,255,0.04);
    backdrop-filter: blur(16px);
    -webkit-backdrop-filter: blur(16px);
    border: 1px solid rgba(255,255,255,0.08);
    border-radius: 16px;
    padding: 16px 20px;
    text-align: center;
    transition: border-color 0.25s, transform 0.25s, box-shadow 0.25s;
    position: relative; overflow: hidden;
}
.lp-stat::before {
    content: '';
    position: absolute; top: 0; left: 0; right: 0; height: 2px;
    border-radius: 16px 16px 0 0;
    opacity: 0.8;
}
.lp-stat-v .n  { color: #C4B5FD !important; }
.lp-stat-v::before { background: linear-gradient(90deg, #8B5CF6, #6366F1); }
.lp-stat-p .n  { color: #F9A8D4 !important; }
.lp-stat-p::before { background: linear-gradient(90deg, #EC4899, #DB2777); }
.lp-stat-c .n  { color: #67E8F9 !important; }
.lp-stat-c::before { background: linear-gradient(90deg, #22D3EE, #0EA5E9); }
.lp-stat-g .n  { color: #6EE7B7 !important; }
.lp-stat-g::before { background: linear-gradient(90deg, #10B981, #059669); }
.lp-stat:hover {
    border-color: rgba(139,92,246,0.28);
    transform: translateY(-3px);
    box-shadow: 0 12px 36px rgba(0,0,0,0.30);
}
.lp-stat .n {
    font-size: 2rem; font-weight: 900;
    letter-spacing: -0.04em; line-height: 1;
    margin-bottom: 5px;
    display: block;
}
.lp-stat .l {
    font-size: 0.58rem; font-weight: 700; letter-spacing: 0.12em;
    text-transform: uppercase; color: #6B7280;
}

/* ══ FAQ / ANSWER STUDIO (below hero) ════════════════════ */
.cat-label {
    display: inline-flex; align-items: center; gap: 7px;
    background: #DBEAFE;
    border: 1px solid #93C5FD;
    color: #1D4ED8 !important; padding: 6px 18px; border-radius: 20px;
    font-size: 0.74rem; font-weight: 700; letter-spacing: 0.05em;
    margin: 26px 0 14px;
}

/* Answer blocks — blue text on very light blue */
.faq-answer-wrap {
    padding: 16px 18px 20px;
    background: #EFF6FF;
    border-radius: 0 0 12px 12px;
}
.faq-answer-label {
    font-size: 0.67rem; font-weight: 700; letter-spacing: 0.14em;
    text-transform: uppercase; color: #1D4ED8; margin-bottom: 10px;
}
.faq-answer-body {
    font-size: 0.91rem; color: #1E3A8A;
    line-height: 1.86;
    padding: 14px 18px;
    background: #DBEAFE;
    border-left: 3px solid #3B82F6;
    border-radius: 0 12px 12px 0;
}
.faq-answer-body strong, .faq-answer-body b { color: #1D4ED8; font-weight: 600; }
.faq-answer-body em { color: #2563EB; font-style: italic; }
.faq-answer-body a { color: #2563EB !important; }

/* ── Expanders — blue on light blue ── */
[data-testid="stExpander"] {
    background: #EFF6FF !important;
    border: 1px solid #BFDBFE !important;
    border-radius: 12px !important;
    overflow: hidden !important;
}
[data-testid="stExpander"] > details > summary {
    background: #EFF6FF !important;
    color: #1E3A8A !important;
    font-size: 0.9rem !important;
    font-weight: 600 !important;
    padding: 13px 18px !important;
}
[data-testid="stExpander"] > details > summary:hover {
    background: #DBEAFE !important;
    color: #1D4ED8 !important;
}
[data-testid="stExpander"] > details > summary p,
[data-testid="stExpander"] > details > summary span {
    color: #1E3A8A !important;
}
[data-testid="stExpander"] details[open] > summary {
    border-bottom: 1px solid #BFDBFE !important;
}
[data-testid="stExpander"] .streamlit-expanderContent {
    background: #EFF6FF !important;
    color: #1E3A8A !important;
}
[data-testid="stExpander"] .streamlit-expanderContent p,
[data-testid="stExpander"] .streamlit-expanderContent li,
[data-testid="stExpander"] .streamlit-expanderContent span {
    color: #1E3A8A !important;
}

/* General descriptive / body text in main content */
[data-testid="stVerticalBlock"] .stMarkdown p {
    color: #1E3A8A !important;
}
[data-testid="stVerticalBlock"] .stMarkdown li {
    color: #1E3A8A !important;
}
/* Tab labels stay light */
[data-testid="stTabs"] [data-testid="stMarkdownContainer"] p { color: inherit !important; }

/* Source badges */
.faq-wa-badge {
    display: inline-flex; align-items: center; gap: 4px;
    background: rgba(16,185,129,0.08);
    border: 1px solid rgba(16,185,129,0.18);
    border-radius: 10px; padding: 2px 10px;
    font-size: 0.65rem; font-weight: 600; color: #34D399;
}
.faq-doc-badge {
    display: inline-flex; align-items: center; gap: 4px;
    background: rgba(99,102,241,0.08);
    border: 1px solid rgba(99,102,241,0.18);
    border-radius: 10px; padding: 2px 10px;
    font-size: 0.65rem; font-weight: 600; color: #818CF8;
}

/* WA citation */
.wa-cite {
    display: inline-flex; align-items: center; gap: 6px;
    background: rgba(16,185,129,0.07);
    border: 1px solid rgba(16,185,129,0.16);
    border-radius: 10px; padding: 5px 13px;
    font-size: 0.72rem; font-weight: 600; color: #34D399;
    margin-top: 12px;
}

/* No-FAQ empty state */
.no-faq {
    text-align: center; padding: 72px 24px;
    background: rgba(17,24,39,0.48);
    backdrop-filter: blur(14px);
    border-radius: 20px;
    border: 1px solid rgba(255,255,255,0.05);
    box-shadow: 0 4px 28px rgba(0,0,0,0.16);
}
.no-faq-icon { font-size: 3rem; margin-bottom: 16px; opacity: 0.7; }
.no-faq h3 { font-size: 1rem; color: #94A3B8; font-weight: 600; margin-bottom: 8px; }
.no-faq p  { font-size: 0.82rem; color: #6B7280; }

/* WA panel */
.wa-panel {
    background: rgba(16,185,129,0.04);
    backdrop-filter: blur(14px);
    border: 1px solid rgba(16,185,129,0.13);
    border-radius: 16px; padding: 18px 22px; margin-bottom: 20px;
    box-shadow: 0 2px 18px rgba(0,0,0,0.10);
}
.wa-panel-title {
    font-size: 0.75rem; font-weight: 700; color: #10B981;
    letter-spacing: 0.09em; text-transform: uppercase; margin-bottom: 14px;
}
.wa-cards-row { display: flex; flex-wrap: wrap; gap: 12px; }
.wa-chat-card {
    background: rgba(17,24,39,0.62);
    backdrop-filter: blur(12px);
    border: 1px solid rgba(16,185,129,0.13);
    border-radius: 14px; padding: 14px 18px; min-width: 240px; flex: 1;
    transition: border-color 0.2s, box-shadow 0.2s;
}
.wa-chat-card:hover {
    border-color: rgba(16,185,129,0.28);
    box-shadow: 0 4px 22px rgba(0,0,0,0.14);
}
.wa-chat-top  { display: flex; align-items: center; gap: 8px; margin-bottom: 7px; }
.wa-chat-icon { font-size: 1.05rem; }
.wa-chat-name { font-size: 0.82rem; font-weight: 600; color: #9CA3AF; flex: 1; }
.wa-chat-badge {
    font-size: 0.65rem; font-weight: 600; color: #10B981;
    background: rgba(16,185,129,0.08);
    border: 1px solid rgba(16,185,129,0.18);
    border-radius: 10px; padding: 2px 9px;
}
.wa-chat-meta { font-size: 0.72rem; color: #9CA3AF; margin-bottom: 4px; }
.wa-chat-qa   { font-size: 0.7rem; color: #10B981; font-weight: 600; margin-top: 5px; }

/* ── Floating chat widget ─────────────────────────────── */
.cf-shell {
    background: rgba(10,13,24,0.98);
    border: 1px solid rgba(139,92,246,0.22);
    border-radius: 20px;
    box-shadow: 0 28px 70px rgba(0,0,0,0.7), 0 0 0 1px rgba(139,92,246,0.08);
    overflow: hidden;
    animation: cfIn 0.26s cubic-bezier(0.2,0.8,0.4,1) both;
    width: 100%;
}
@keyframes cfIn {
    from { opacity:0; transform:translateY(18px) scale(0.96); }
    to   { opacity:1; transform:translateY(0)     scale(1);   }
}
.cf-header {
    display:flex; align-items:center; justify-content:space-between;
    padding:13px 16px;
    background:linear-gradient(135deg,rgba(17,24,39,1) 0%,rgba(22,30,50,1) 100%);
    border-bottom:1px solid rgba(255,255,255,0.05);
}
.cf-title-row { display:flex; align-items:center; gap:8px; }
.cf-live-dot { width:8px; height:8px; border-radius:50%; background:#10B981;
    box-shadow:0 0 6px rgba(16,185,129,0.7); animation:livepulse 3s ease-in-out infinite; }
.cf-title { font-size:0.87rem; font-weight:700; color:#E5E7EB; }
.cf-sub   { font-size:0.68rem; color:#6B7280; margin-top:1px; }
.cf-msgs {
    padding:12px 14px;
    max-height:320px; min-height:60px;
    overflow-y:auto; display:flex; flex-direction:column; gap:9px;
    scroll-behavior:smooth;
}
.cf-msgs::-webkit-scrollbar { width:3px; }
.cf-msgs::-webkit-scrollbar-thumb { background:rgba(139,92,246,0.35); border-radius:2px; }
.cf-user-bubble {
    align-self:flex-end;
    background:linear-gradient(135deg,#7C3AED,#6D28D9);
    color:#F3F4F6; font-size:0.82rem; line-height:1.5;
    padding:9px 13px; border-radius:16px 16px 4px 16px;
    max-width:88%; word-break:break-word;
}
.cf-ai-bubble {
    align-self:flex-start;
    background:rgba(26,34,52,0.9); color:#D1D5DB;
    font-size:0.82rem; line-height:1.6;
    padding:9px 13px; border-radius:16px 16px 16px 4px;
    max-width:92%; border:1px solid rgba(255,255,255,0.06);
    word-break:break-word;
}
.cf-empty { text-align:center; padding:28px 14px; color:#6B7280; font-size:0.8rem; }
.cf-empty-icon { font-size:1.7rem; margin-bottom:8px; opacity:0.55; display:block; }
/* FAB — Streamlit button override when sitting in the float container */
.cf-fab-wrap button {
    width:58px !important; height:58px !important;
    border-radius:50% !important; padding:0 !important;
    background:linear-gradient(135deg,#8B5CF6,#EC4899) !important;
    color:#fff !important; font-size:1.45rem !important;
    border:none !important;
    box-shadow:0 4px 28px rgba(139,92,246,0.55) !important;
    transition:transform 0.22s cubic-bezier(.34,1.56,.64,1) !important;
    animation:fab-pulse 3s ease-in-out infinite !important;
}
.cf-fab-wrap button:hover { transform:scale(1.1) !important; }
/* Close button in chat header */
.cf-close-btn button {
    padding:2px 10px !important; font-size:0.75rem !important;
    border-radius:8px !important; min-height:28px !important;
    background:rgba(255,255,255,0.06) !important;
    color:#9CA3AF !important; border:1px solid rgba(255,255,255,0.1) !important;
}
.cf-close-btn button:hover { background:rgba(255,255,255,0.12) !important; color:#E5E7EB !important; }
/* Input row inside float panel */
.cf-input-row { padding:8px 12px 10px; border-top:1px solid rgba(255,255,255,0.05); }
.cf-input-row .stTextInput > div > div > input {
    background:rgba(17,24,39,0.9) !important; border:1px solid rgba(99,102,241,0.25) !important;
    border-radius:10px !important; color:#E5E7EB !important; font-size:0.83rem !important;
}
.cf-input-row .stTextInput > div > div > input:focus {
    border-color:rgba(139,92,246,0.55) !important; box-shadow:0 0 0 2px rgba(139,92,246,0.15) !important;
}

/* Nav pill buttons */
div[data-testid="stHorizontalBlock"]:has(button[key="faq_nav_btn"]) .stButton > button,
div[data-testid="stHorizontalBlock"]:has(button[key="settings_btn"]) .stButton > button,
div[data-testid="stHorizontalBlock"]:has(button[key="chat_nav_btn"]) .stButton > button,
div[data-testid="stHorizontalBlock"]:has(button[key="settings_btn_faq"]) .stButton > button {
    border-radius: 20px !important;
}

@keyframes fab-pulse {
    0%,100% { box-shadow:0 4px 28px rgba(139,92,246,0.55),0 0 0 0 rgba(139,92,246,0.25); }
    50%      { box-shadow:0 4px 28px rgba(139,92,246,0.55),0 0 0 10px rgba(139,92,246,0); }
}
</style>
""", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════
#  STORAGE  —  Streamlit-native JSON file (kb_store.json)
# ══════════════════════════════════════════════════════════════════
def load_kb():
    if st.session_state.get("_kb_loaded"):
        return
    if os.path.exists(KB_FILE):
        try:
            with open(KB_FILE, "r", encoding="utf-8") as f:
                d = json.load(f)
            for k in KB_KEYS:
                st.session_state[k] = d.get(k, [])
            st.session_state["show_sources"] = d.get("show_sources", False)
        except Exception:
            pass
    st.session_state["_kb_loaded"] = True

def save_kb():
    data = {k: st.session_state.get(k, []) for k in KB_KEYS}
    data["show_sources"] = st.session_state.get("show_sources", False)
    with open(KB_FILE, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

def kb_stats():
    docs  = len(st.session_state.get("kb_documents", []))
    links = len(st.session_state.get("kb_links",     []))
    wa    = len(st.session_state.get("kb_whatsapp",  []))
    pages = len(st.session_state.get("kb_crawled",   []))
    return docs, links, wa, pages

def total_sources():
    return sum(kb_stats())


# ══════════════════════════════════════════════════════════════════
#  SESSION STATE
# ══════════════════════════════════════════════════════════════════
_DEFAULTS = {
    "page":          "faq",
    "chat_open":     False,
    "chat_history":  [],
    "_last_input":   "",
    "_mini_last":    "",
    "quick_q":       "",
    "show_sources":  False,
    "_kb_loaded":    False,
    "kb_faqs":       [],
    **{k: [] for k in KB_KEYS},
}
for k, v in _DEFAULTS.items():
    if k not in st.session_state:
        st.session_state[k] = v

load_kb()

# Handle FAB / link navigation via query params
if "nav" in st.query_params:
    _nav_target = st.query_params["nav"]
    st.query_params.clear()
    if _nav_target == "chat":
        st.session_state.chat_open = True   # open panel, don't leave the page
    elif _nav_target in ("faq", "settings"):
        st.session_state.page = _nav_target


# ══════════════════════════════════════════════════════════════════
#  FILE PARSERS
# ══════════════════════════════════════════════════════════════════
def _ext(filename): return filename.lower().rsplit(".", 1)[-1] if "." in filename else ""

def FILE_ICON(ext):
    return {
        "pdf":"📄","txt":"📝","md":"📝","docx":"📋","doc":"📋","odt":"📋",
        "xlsx":"📊","xls":"📊","xlsm":"📊","ods":"📊","csv":"📊",
        "pptx":"🖼️","ppt":"🖼️","odp":"🖼️",
        "html":"🌐","htm":"🌐","mhtml":"🌐","xhtml":"🌐",
        "json":"🗂️","jsonl":"🗂️","yaml":"🗂️","xml":"🗂️",
        "epub":"📚","rtf":"📄","eml":"📧","msg":"📧",
        "ipynb":"📓","py":"🐍","js":"📜","ts":"📜","sql":"🗄️",
        "sh":"⚙️","rb":"💎","go":"🐹","rs":"⚙️","java":"☕",
    }.get(ext, "📎")

def parse_file(f) -> str:
    import io
    name = f.name.lower(); raw = f.read(); ext = _ext(name)

    # ── Plain-text formats (direct decode) ───────────────────────
    PLAIN_TEXT = {
        "txt","md","markdown","rst","log","csv","tsv",
        "yaml","yml","toml","ini","cfg","env",
        "sh","bash","zsh","fish",
        "py","js","ts","jsx","tsx","mjs","cjs",
        "json","jsonl","ndjson",
        "sql","r","rb","php","java","cpp","c","h","hpp",
        "go","rs","swift","kt","cs","lua","tex","scala","pl",
        "xml","svg","rss","atom",
    }
    if ext in PLAIN_TEXT:
        return raw.decode("utf-8", errors="ignore")

    # ── HTML / HTM — strip tags, keep readable text ───────────────
    if ext in ("html","htm","mhtml","mht","xhtml"):
        try:
            from bs4 import BeautifulSoup as BS
            soup = BS(raw.decode("utf-8","ignore"), "html.parser")
            for tag in soup(["script","style","nav","footer","header",
                             "aside","noscript","meta","link"]):
                tag.decompose()
            title = soup.title.string.strip() if soup.title else ""
            text  = soup.get_text("\n", strip=True)
            return (f"[Title: {title}]\n\n" if title else "") + text
        except Exception as e:
            return raw.decode("utf-8","ignore")

    # ── PDF ───────────────────────────────────────────────────────
    if ext == "pdf":
        try:
            import pdfplumber
            with pdfplumber.open(io.BytesIO(raw)) as pdf:
                return "\n\n".join(p.extract_text() or "" for p in pdf.pages)
        except Exception:
            try:
                import PyPDF2
                r = PyPDF2.PdfReader(io.BytesIO(raw))
                return "\n\n".join(p.extract_text() or "" for p in r.pages)
            except Exception as e: return f"[PDF error: {e}]"

    # ── Word ──────────────────────────────────────────────────────
    if ext in ("docx","doc","odt"):
        try:
            import docx
            d = docx.Document(io.BytesIO(raw))
            parts = [p.text for p in d.paragraphs if p.text.strip()]
            for tbl in d.tables:
                for row in tbl.rows:
                    parts.append(" | ".join(c.text for c in row.cells))
            return "\n".join(parts)
        except Exception as e: return f"[DOCX error: {e}]"

    # ── Excel ─────────────────────────────────────────────────────
    if ext in ("xlsx","xls","xlsm","ods"):
        try:
            import openpyxl
            wb = openpyxl.load_workbook(io.BytesIO(raw), read_only=True, data_only=True)
            parts = []
            for sn in wb.sheetnames:
                ws = wb[sn]; parts.append(f"=== Sheet: {sn} ===")
                for row in ws.iter_rows(values_only=True):
                    s = " | ".join(str(c) for c in row if c is not None)
                    if s.strip(): parts.append(s)
            return "\n".join(parts)
        except Exception as e: return f"[Excel error: {e}]"

    # ── PowerPoint ────────────────────────────────────────────────
    if ext in ("pptx","ppt","odp"):
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(raw)); parts = []
            for i,sl in enumerate(prs.slides,1):
                parts.append(f"=== Slide {i} ===")
                for sh in sl.shapes:
                    if hasattr(sh,"text") and sh.text.strip(): parts.append(sh.text)
            return "\n".join(parts)
        except Exception as e: return f"[PPTX error: {e}]"

    # ── RTF ───────────────────────────────────────────────────────
    if ext == "rtf":
        try:
            from striprtf.striprtf import rtf_to_text
            return rtf_to_text(raw.decode("utf-8","ignore"))
        except Exception as e: return f"[RTF error: {e}]"

    # ── EPUB ──────────────────────────────────────────────────────
    if ext == "epub":
        try:
            import ebooklib; from ebooklib import epub; from bs4 import BeautifulSoup as BS
            book = epub.read_epub(io.BytesIO(raw)); parts = []
            for item in book.get_items_of_type(ebooklib.ITEM_DOCUMENT):
                soup = BS(item.get_content(),"html.parser")
                t = soup.get_text(" ",strip=True)
                if t: parts.append(t)
            return "\n\n".join(parts)
        except Exception as e: return f"[EPUB error: {e}]"

    # ── Email (.eml / .msg) ───────────────────────────────────────
    if ext in ("eml","msg"):
        try:
            import email
            msg = email.message_from_bytes(raw)
            parts = []
            for hdr in ("From","To","Subject","Date"):
                if msg.get(hdr): parts.append(f"{hdr}: {msg[hdr]}")
            parts.append("")
            for part in msg.walk():
                ct = part.get_content_type()
                if ct == "text/plain":
                    parts.append(part.get_payload(decode=True).decode("utf-8","ignore"))
                elif ct == "text/html":
                    try:
                        from bs4 import BeautifulSoup as BS
                        parts.append(BS(part.get_payload(decode=True),"html.parser").get_text("\n",strip=True))
                    except Exception:
                        pass
            return "\n".join(parts)
        except Exception as e: return f"[EML error: {e}]"

    # ── Notebook (.ipynb) ─────────────────────────────────────────
    if ext == "ipynb":
        try:
            nb = json.loads(raw.decode("utf-8","ignore"))
            parts = []
            for cell in nb.get("cells",[]):
                ct = cell.get("cell_type","")
                src = "".join(cell.get("source",[]))
                if ct == "markdown":
                    parts.append(src)
                elif ct == "code":
                    parts.append(f"```python\n{src}\n```")
                for out in cell.get("outputs",[]):
                    text = "".join(out.get("text",""))
                    if text: parts.append(text)
            return "\n\n".join(parts)
        except Exception as e: return f"[IPYNB error: {e}]"

    # ── Fallback: try raw text decode ─────────────────────────────
    for enc in ("utf-8","latin-1","cp1252"):
        try: return raw.decode(enc)
        except: continue
    return f"[Unsupported format: {f.name}]"


# ══════════════════════════════════════════════════════════════════
#  WEB HELPERS
# ══════════════════════════════════════════════════════════════════
def _http(url):
    import requests
    return requests.get(url, headers={"User-Agent":"Mozilla/5.0 (ConvinBot/1.0)"}, timeout=12)

def _scrape(soup, url):
    for t in soup(["script","style","nav","footer","header","aside","noscript"]):
        t.decompose()
    title = (soup.title.string.strip() if soup.title else url)[:120]
    text  = "\n".join(
        p.get_text(" ",strip=True)
        for p in soup.find_all(["p","h1","h2","h3","h4","li","td","blockquote"])
        if len(p.get_text(strip=True)) > 25
    )
    return title, text[:16000]

def fetch_url(url):
    try:
        from bs4 import BeautifulSoup
        r = _http(url); r.raise_for_status()
        return _scrape(BeautifulSoup(r.text,"html.parser"), url)
    except Exception as e:
        return url, f"[Error: {e}]"

_WA_SKIP = {
    "<Media omitted>", "This message was deleted", "image omitted",
    "video omitted", "audio omitted", "sticker omitted", "document omitted",
    "GIF omitted", "Contact card omitted", "Missed voice call",
    "Missed video call", "null", "",
}

# WhatsApp message patterns — iOS and Android, 12h and 24h
_WA_PATTERNS = [
    # iOS:     [DD/MM/YYYY, HH:MM:SS] Sender: Msg
    re.compile(r"^\[(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4}),\s*(\d{1,2}:\d{2}(?::\d{2})?(?:\s*[AaPp][Mm])?)\]\s+([^:]+):\s*(.*)"),
    # Android: DD/MM/YYYY, HH:MM - Sender: Msg
    re.compile(r"^(\d{1,2}[\/\-\.]\d{1,2}[\/\-\.]\d{2,4}),?\s+(\d{1,2}:\d{2}(?::\d{2})?(?:\s*[AaPp][Mm])?)\s*[-–]\s*([^:]+):\s*(.*)"),
]

def _wa_match(line: str):
    """Try all WA patterns; return (date, time, sender, msg) or None."""
    for pat in _WA_PATTERNS:
        m = pat.match(line.strip())
        if m:
            return m.group(1), m.group(2).strip(), m.group(3).strip(), m.group(4).strip()
    return None

def parse_wa(raw: str) -> str:
    """Parse a WhatsApp export preserving full [date time] Sender: message format."""
    lines_out = []
    for ln in raw.split("\n"):
        s = ln.strip()
        if not s:
            continue
        hit = _wa_match(s)
        if hit:
            date, time_, sender, msg = hit
            if not msg or msg in _WA_SKIP or any(sk.lower() in msg.lower() for sk in _WA_SKIP):
                continue
            lines_out.append(f"[{date} {time_}] {sender}: {msg}")
        elif lines_out:
            lines_out[-1] = lines_out[-1] + " " + s  # continuation
    return "\n".join(lines_out)

def parse_wa_meta(content: str) -> dict:
    """Extract chat metadata for the Answer Studio pre-analysis header."""
    messages = []
    for ln in content.split("\n"):
        hit = _wa_match(ln)
        if hit:
            date, time_, sender, msg = hit
            messages.append({"date": date, "time": time_, "sender": sender, "text": msg})

    if not messages:
        return {"valid": False, "total": 0, "participants": [], "date_range": ""}

    participants = list(dict.fromkeys(m["sender"] for m in messages))
    counts = {}
    for m in messages:
        counts[m["sender"]] = counts.get(m["sender"], 0) + 1

    return {
        "valid": True,
        "total": len(messages),
        "participants": participants,
        "msg_counts": counts,
        "first_date": messages[0]["date"],
        "last_date": messages[-1]["date"],
        "date_range": f"{messages[0]['date']} → {messages[-1]['date']}",
    }

def crawl_site(root, max_p, status_ph, prog_ph):
    try:
        import requests; from bs4 import BeautifulSoup
    except ImportError:
        return 0
    pr = urlparse(root); base = f"{pr.scheme}://{pr.netloc}"; host = pr.netloc
    SKIP = {".pdf",".png",".jpg",".jpeg",".gif",".svg",".zip",".exe",
            ".mp4",".mp3",".webp",".ico",".css",".js",".woff",".ttf"}
    def ok(h):
        p = urlparse(h)
        return (not p.netloc or p.netloc==host) and \
               not any(p.path.lower().endswith(e) for e in SKIP) and \
               p.scheme in ("http","https","")
    visited = set(); queue = [root]; done = 0
    existing = {p["url"] for p in st.session_state.kb_crawled}; new_pages = []
    while queue and done < max_p:
        url = queue.pop(0).split("#")[0].rstrip("/") or root
        if url in visited: continue
        visited.add(url)
        try:
            resp = _http(url)
            if "text/html" not in resp.headers.get("Content-Type",""): continue
            soup = BeautifulSoup(resp.text,"html.parser")
            for a in soup.find_all("a",href=True):
                h = urljoin(base,a["href"]).split("#")[0].rstrip("/")
                if ok(h) and h not in visited: queue.append(h)
            title, text = _scrape(soup, url)
            if not text.strip(): continue
            done += 1
            pg = {"url":url,"title":title,"content":text,
                  "added_at":datetime.now().isoformat(),"size":len(text)}
            if url not in existing:
                new_pages.append(pg); existing.add(url)
            prog_ph.progress(min(done/max_p,1.0))
            status_ph.caption(f"🕷️  {done}/{max_p} — {url[:60]}")
        except Exception: continue
    st.session_state.kb_crawled.extend(new_pages)
    prog_ph.empty(); status_ph.empty()
    return done


# ══════════════════════════════════════════════════════════════════
#  KNOWLEDGE BASE CONTEXT
# ══════════════════════════════════════════════════════════════════
def build_context():
    """Return (full_context_string, list_of_source_names)."""
    parts, names = [], []
    for d in st.session_state.kb_documents:
        parts.append(f"=== DOCUMENT: {d['name']} ===\n{d['content']}")
        names.append(d["name"])
    for l in st.session_state.kb_links:
        parts.append(f"=== WEBSITE: {l['title']} ===\n{l['content']}")
        names.append(l["title"])
    for w in st.session_state.kb_whatsapp:
        parts.append(f"=== WHATSAPP CHAT: {w['name']} ===\n{w['content']}")
        names.append(w["name"])
    for p in st.session_state.kb_crawled:
        parts.append(f"=== PAGE: {p['title']} ===\n{p['content']}")
        names.append(p["title"])
    ctx = "\n\n".join(parts)
    return ctx[:MAX_CTX], names


def build_chat_context() -> tuple[str, list[str]]:
    """Compact FAQ-based context for chat — much faster than raw pages."""
    parts, names = [], []
    faqs = st.session_state.get("kb_faqs", [])
    if faqs:
        qa_lines = [f"Q: {f['question']}\nA: {f['answer']}" for f in faqs]
        parts.append("=== KNOWLEDGE BASE ===\n" + "\n\n".join(qa_lines))
        names.append(f"{len(faqs)} Q&A pairs")
    for d in st.session_state.get("kb_documents", []):
        parts.append(f"=== DOCUMENT: {d['name']} ===\n{d['content'][:6000]}")
        names.append(d["name"])
    ctx = "\n\n".join(parts)
    return ctx[:180_000], names


# ══════════════════════════════════════════════════════════════════
#  CLAUDE  —  full context, prompt caching
# ══════════════════════════════════════════════════════════════════
def get_client():
    key = st.secrets.get("ANTHROPIC_API_KEY", os.getenv("ANTHROPIC_API_KEY",""))
    if not key:
        raise ValueError("ANTHROPIC_API_KEY not configured.")
    return anthropic.Anthropic(api_key=key)

def ask_claude(query: str) -> tuple[str, list[str]]:
    ctx, names = build_context()
    if not ctx.strip():
        return (
            "I don't have any knowledge base loaded yet. "
            "Please ask your admin to add documents or resources.",
            [],
        )

    client = get_client()

    SYSTEM_INSTRUCTIONS = (
        "You are a professional Customer Support AI for Convin.\n\n"
        "Behavior:\n"
        "• Answer ONLY from the knowledge base provided — never fabricate facts.\n"
        "• Read the ENTIRE knowledge base before responding, including WhatsApp chats.\n"
        "• Structure every response:\n"
        "    1. Direct answer (1–2 sentences)\n"
        "    2. Steps or details (numbered list if 3+ steps, otherwise inline)\n"
        "    3. Brief explanation (only if it adds value)\n"
        "    4. Friendly closing line\n"
        "• Use **bold** for key terms, product names, and important values.\n"
        "• Keep answers concise — no filler words, no repetition.\n"
        "• If the answer is NOT in the knowledge base, respond with:\n"
        "  \"This might need further verification — let me connect you with "
        "the right person from our team.\"\n"
        "• For WhatsApp conversations: ALWAYS cite the source with this exact format:\n"
        "  > 💬 Chatted by [Sender Name] at [Time] on [Date]\n"
        "  Include the actual sender name, time, and date from the message metadata.\n"
        "• Do NOT say 'based on the document' or 'according to the file' for documents.\n"
        "• Do NOT reveal file names or document titles in answers.\n"
        "• WhatsApp references are valuable citations — always include them when used.\n"
    )

    system = [
        {"type": "text", "text": SYSTEM_INSTRUCTIONS},
        {"type": "text", "text": "KNOWLEDGE BASE:\n" + ctx,
         "cache_control": {"type": "ephemeral"}},
    ]

    history = [{"role": m["role"], "content": m["content"]}
               for m in st.session_state.chat_history[-10:]]
    history.append({"role": "user", "content": query})

    r = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=900,
        system=system,
        messages=history,
        extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"},
    )
    return r.content[0].text, names


def ask_claude_stream(query: str, placeholder) -> tuple[str, list[str]]:
    """Stream Claude response token-by-token into a Streamlit placeholder."""
    ctx, names = build_chat_context()
    if not ctx.strip():
        msg = ("I don't have any knowledge base loaded yet. "
               "Please ask your admin to add documents or resources.")
        placeholder.markdown(msg)
        return msg, []

    SYSTEM = (
        "You are Animesh, a professional support AI for Convin Sense.\n"
        "• Answer ONLY from the knowledge base — never fabricate facts.\n"
        "• Be concise and direct. Use **bold** for key terms and product names.\n"
        "• If the answer is not in the knowledge base respond with: "
        "\"This might need further verification — let me connect you with the right person from our team.\"\n"
        "• For WhatsApp citations use: > 💬 [Name] on [Date]\n"
    )
    system = [
        {"type": "text", "text": SYSTEM},
        {"type": "text", "text": "KNOWLEDGE BASE:\n" + ctx,
         "cache_control": {"type": "ephemeral"}},
    ]
    history = [{"role": m["role"], "content": m["content"]}
               for m in st.session_state.chat_history[-6:]]
    history.append({"role": "user", "content": query})

    full_text = ""
    try:
        with get_client().messages.stream(
            model="claude-haiku-4-5-20251001",
            max_tokens=600,
            system=system,
            messages=history,
            extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"},
        ) as stream:
            for text in stream.text_stream:
                full_text += text
                placeholder.markdown(full_text + " ▌")
        placeholder.markdown(full_text)
    except Exception as e:
        full_text = f"Something went wrong — please try again. *(Error: {e})*"
        placeholder.markdown(full_text)

    return full_text, names


# ══════════════════════════════════════════════════════════════════
#  FAQ GENERATOR
# ══════════════════════════════════════════════════════════════════
def _faq_call(client, prompt: str) -> list[dict]:
    """Single Claude call → parsed FAQ list."""
    r = client.messages.create(
        model="claude-sonnet-4-6",
        max_tokens=16000,
        messages=[{
            "role": "user",
            "content": [{"type": "text", "text": prompt,
                         "cache_control": {"type": "ephemeral"}}],
        }],
        extra_headers={"anthropic-beta": "prompt-caching-2024-07-31"},
    )
    raw = re.sub(r"^```[a-z]*\n?", "", r.content[0].text.strip())
    raw = re.sub(r"\n?```$", "", raw)
    try:
        items = json.loads(raw)
    except json.JSONDecodeError:
        pairs = re.findall(
            r'"question"\s*:\s*"([^"]+)"[^}]*"answer"\s*:\s*"([^"]+)"', raw
        )
        return [{"category": "General", "question": q, "answer": a} for q, a in pairs]
    return [
        {
            "category": str(i.get("category", "General")),
            "question": str(i.get("question", i.get("q", ""))),
            "answer":   str(i.get("answer",   i.get("a", ""))),
        }
        for i in items if isinstance(i, dict) and i.get("question")
    ]


def generate_faqs(progress_cb=None) -> list[dict]:
    """Multi-pass FAQ generation — one full Claude call per source type.

    progress_cb(pct, label) is called after each pass so the UI can update.
    """
    client  = get_client()
    all_faqs: list[dict] = []

    BASE_RULES = (
        "Rules:\n"
        "• Only include facts present in the content — never invent.\n"
        "• Be EXHAUSTIVE: cover every topic, feature, process, policy, edge case, "
        "decision, question, issue, and fact — no matter how minor.\n"
        "• Each answer: 2–6 sentences or a short bullet list.\n"
        "• Group into specific, descriptive categories.\n"
        "• Return ONLY a raw JSON array (no markdown, no extra text):\n"
        '[\n  {"category":"Name","question":"Q?","answer":"A."},\n  ...\n]\n'
    )

    # ── Pass 1: Documents ─────────────────────────────────────────
    docs = st.session_state.get("kb_documents", [])
    if docs:
        ctx = "\n\n".join(
            f"=== DOCUMENT: {d['name']} ===\n{d['content']}" for d in docs
        )[:MAX_CTX]
        prompt = (
            "You are extracting the maximum possible FAQs from uploaded documents.\n"
            "Read every sentence. Extract every question a user could ever ask.\n\n"
            + BASE_RULES +
            "\nExtract the MAXIMUM number of Q&A pairs possible.\n\n"
            "DOCUMENTS:\n" + ctx
        )
        if progress_cb: progress_cb(0.15, f"📄 Processing {len(docs)} document(s)…")
        try:
            all_faqs.extend(_faq_call(client, prompt))
        except Exception:
            pass

    # ── Pass 2: Web links + crawled pages ────────────────────────
    web = st.session_state.get("kb_links", []) + st.session_state.get("kb_crawled", [])
    if web:
        ctx = "\n\n".join(
            f"=== PAGE: {p.get('title', p.get('url',''))} ===\n{p['content']}"
            for p in web
        )[:MAX_CTX]
        prompt = (
            f"You are extracting the maximum possible FAQs from {len(web)} web page(s).\n"
            "Read every sentence. Extract every question a user could ever ask.\n\n"
            + BASE_RULES +
            "\nExtract the MAXIMUM number of Q&A pairs possible.\n\n"
            "WEB PAGES:\n" + ctx
        )
        if progress_cb: progress_cb(0.45, f"🌐 Processing {len(web)} web page(s)…")
        try:
            all_faqs.extend(_faq_call(client, prompt))
        except Exception:
            pass

    # ── Pass 3 & 4: WhatsApp — two deep-extraction passes per chat ──
    wa_chats = st.session_state.get("kb_whatsapp", [])
    if wa_chats:
        total_wa = len(wa_chats)
        for wi, chat in enumerate(wa_chats):
            content = chat.get("content", "").strip()
            if not content:
                continue

            meta = parse_wa_meta(content)
            plist = ", ".join(meta["participants"]) if meta["valid"] else "unknown"
            drange = meta.get("date_range", "")
            total_msg = meta.get("total", "?")

            pct_a = 0.60 + (wi / total_wa) * 0.15
            pct_b = pct_a + 0.07

            CITE_RULE = (
                "CITATION RULE — every answer MUST end with one of:\n"
                "  • Single person: '💬 [Name] on [Date] at [Time]'\n"
                "  • Exchange: '💬 [Name1] asked, [Name2] replied on [Date]'\n"
                "  • Group: '💬 Discussed by [Name1], [Name2] on [Date]'\n"
                "Use the actual names, dates, and times from the messages.\n\n"
            )

            HEADER = (
                f"WhatsApp chat: {chat['name']}\n"
                f"Participants: {plist}\n"
                f"Date range: {drange}  |  Messages: {total_msg}\n\n"
                "Each line format: [DD/MM/YY HH:MM] Sender: Message\n\n"
            )

            # ── Sub-pass A: Direct Q&As + Decisions + Action items ──
            if progress_cb:
                progress_cb(pct_a, f"💬 Chat {wi+1}/{total_wa} — extracting Q&As, decisions, actions…")
            prompt_a = (
                "You are a knowledge analyst. Deeply read this WhatsApp conversation.\n\n"
                + HEADER
                + CITE_RULE
                + "EXTRACT THESE THREE TYPES (be exhaustive):\n\n"
                "TYPE 1 — QUESTIONS & ANSWERS\n"
                "Find every explicit question (ends with ?, starts with how/what/when/why/who/where/can/should/is/are/do/does/will/would) "
                "AND every implied question (topic raised that others responded to). "
                "Pair each with its answer from the conversation.\n"
                "→ Category: 'WhatsApp: Questions & Answers'\n\n"
                "TYPE 2 — DECISIONS & AGREEMENTS\n"
                "Find everything agreed upon, confirmed, resolved, or decided. "
                "Look for: 'agreed', 'decided', 'confirmed', 'let's go with', 'we'll', 'done', 'sorted', 'ok let's'.\n"
                "→ Category: 'WhatsApp: Decisions & Agreements'\n\n"
                "TYPE 3 — ACTION ITEMS & TASKS\n"
                "Find every task assigned, promise made, or next step defined. "
                "Look for: 'will do', 'I'll', 'you need to', 'please', 'can you', 'by [date]', 'follow up'.\n"
                "→ Category: 'WhatsApp: Action Items & Tasks'\n\n"
                "Return ONLY raw JSON array:\n"
                '[\n  {"category":"WhatsApp: Questions & Answers","question":"Q?","answer":"A. 💬 ..."},\n  ...\n]\n\n'
                "Be exhaustive — extract every single instance.\n\n"
                "CHAT:\n" + content[:MAX_CTX]
            )
            try:
                all_faqs.extend(_faq_call(client, prompt_a))
            except Exception:
                pass

            # ── Sub-pass B: Information + Problems + Context + Insights ──
            if progress_cb:
                progress_cb(pct_b, f"💬 Chat {wi+1}/{total_wa} — extracting knowledge, issues, insights…")
            prompt_b = (
                "You are a knowledge analyst. Deeply read this WhatsApp conversation.\n\n"
                + HEADER
                + CITE_RULE
                + "EXTRACT THESE FOUR TYPES (be exhaustive):\n\n"
                "TYPE 4 — INFORMATION & KNOWLEDGE SHARED\n"
                "Find every fact, figure, process, instruction, contact, link, or data point shared. "
                "Turn each into 'What is/How does/What was...?' format.\n"
                "→ Category: 'WhatsApp: Information & Knowledge'\n\n"
                "TYPE 5 — PROBLEMS & RESOLUTIONS\n"
                "Find every issue, complaint, bug, confusion, or blocker raised — and how it was resolved. "
                "If unresolved, note that.\n"
                "→ Category: 'WhatsApp: Issues & Resolutions'\n\n"
                "TYPE 6 — BUSINESS & PRODUCT INSIGHTS\n"
                "Find anything about clients, products, features, pricing, timelines, deals, or strategy.\n"
                "→ Category: 'WhatsApp: Business & Product Insights'\n\n"
                "TYPE 7 — CONTEXT & BACKGROUND\n"
                "Find any background context, history, or explanations given about the situation or topic.\n"
                "→ Category: 'WhatsApp: Context & Background'\n\n"
                "Return ONLY raw JSON array:\n"
                '[\n  {"category":"WhatsApp: Information & Knowledge","question":"Q?","answer":"A. 💬 ..."},\n  ...\n]\n\n'
                "Be exhaustive — extract every single piece of knowledge.\n\n"
                "CHAT:\n" + content[:MAX_CTX]
            )
            try:
                all_faqs.extend(_faq_call(client, prompt_b))
            except Exception:
                pass

    if progress_cb: progress_cb(0.95, "✨ Deduplicating and finalising…")

    # ── Deduplicate by question (case-insensitive) ────────────────
    seen: set[str] = set()
    deduped: list[dict] = []
    for f in all_faqs:
        key_q = f["question"].lower().strip()
        if key_q and key_q not in seen:
            seen.add(key_q)
            deduped.append(f)

    return deduped


# ══════════════════════════════════════════════════════════════════
#  SHARED TOP NAV
# ══════════════════════════════════════════════════════════════════
def render_topnav(show_settings_btn=True, show_back_btn=False, show_chat_btn=False):
    docs, links, wa, pages = kb_stats()
    total = docs + links + wa + pages
    status_label = f"{total} sources loaded" if total else "No knowledge base"

    # HTML-only portion (purely visual)
    logo_html = (
        f'<img src="{_LOGO_URI}" class="topnav-logo" alt="Convin">'
        if _LOGO_URI else
        '<div class="dot">K</div>'
    )
    st.markdown(f"""
    <div class="topnav">
      <div class="topnav-brand">
        {logo_html}
        <span class="name">Convin Klaro</span>
        <span class="badge">AI Support Intelligence</span>
      </div>
      <div class="topnav-right">
        <span class="topnav-status">
          <span class="live-dot"></span>
          {status_label}
        </span>
      </div>
    </div>
    """, unsafe_allow_html=True)

    # Streamlit buttons — functional nav row
    nav_spacer = st.container()
    with nav_spacer:
        c_l, c_faq, c_set = st.columns([7, 1, 1])
        if show_back_btn:
            with c_l:
                if st.button("← Back to Answer Studio", key="back_btn", type="secondary"):
                    st.session_state.page = "faq"
                    st.rerun()
        if show_settings_btn:
            with c_faq:
                if st.button("✦ Answer Studio", key="faq_nav_btn", type="secondary",
                             use_container_width=True):
                    st.session_state.page = "faq"
                    st.rerun()
            with c_set:
                if st.button("⚙ Settings", key="settings_btn", type="secondary",
                             use_container_width=True):
                    st.session_state.page = "settings"
                    st.rerun()
        if show_chat_btn:
            chat_open = st.session_state.get("chat_open", False)
            with c_faq:
                btn_label = "✕ Chat" if chat_open else "💬 Chat"
                if st.button(btn_label, key="chat_nav_btn", type="secondary",
                             use_container_width=True):
                    st.session_state.chat_open = not chat_open
                    st.rerun()
            with c_set:
                if st.button("⚙ Settings", key="settings_btn_faq", type="secondary",
                             use_container_width=True):
                    st.session_state.page = "settings"
                    st.rerun()


# ══════════════════════════════════════════════════════════════════
#  CHAT PAGE
# ══════════════════════════════════════════════════════════════════
SUGGESTIONS = [
    "🤖  What is Convin AI?",
    "📊  How does call scoring work?",
    "🔗  What integrations are supported?",
    "⚡  How to set up Auto QA?",
    "💬  Any recent discussions or decisions?",
    "📤  How to export reports?",
]

def render_chat():
    render_topnav(show_settings_btn=True, show_back_btn=False)

    # ── Center column for chat
    _, main, _ = st.columns([1, 5, 1])
    with main:

        # ── Message feed
        if not st.session_state.chat_history:
            # Welcome screen — live KB stats
            docs, links, wa, pages = kb_stats()
            total = docs + links + wa + pages
            stat_chips = []
            if docs:
                stat_chips.append(f'<div class="kb-stat-chip active"><span class="num">{docs}</span><span>📄 Docs</span></div>')
            if links + pages:
                stat_chips.append(f'<div class="kb-stat-chip active"><span class="num">{links+pages}</span><span>🌐 Web pages</span></div>')
            if wa:
                stat_chips.append(f'<div class="kb-stat-chip active wa"><span class="num">{wa}</span><span>💬 WA chats</span></div>')

            if stat_chips:
                stats_html = '<div class="kb-stats-row">' + "".join(stat_chips) + '</div>'
                ready_badge = '<div class="ready-badge">✦ Knowledge base ready</div>'
            else:
                stats_html = '<div class="no-kb-hint">No sources loaded yet — add documents in ⚙ Settings</div>'
                ready_badge = ''

            st.markdown(f"""
            <div class="welcome-card">
              <div class="welcome-icon">✦</div>
              <div class="welcome-eyebrow">Convin Klaro</div>
              <div class="welcome-title">Hi, I'm Animesh. <span>How may I help you?</span></div>
              <div class="welcome-sub">
                Ask me anything — I'll search across all your docs, web pages &amp; chats to find the answer.
              </div>
              {stats_html}
              {ready_badge}
            </div>
            """, unsafe_allow_html=True)

            # Suggestion chips — two rows of 3
            row1, row2 = SUGGESTIONS[:3], SUGGESTIONS[3:]
            scols1 = st.columns(3)
            for i, (col, q) in enumerate(zip(scols1, row1)):
                with col:
                    if st.button(q, key=f"sugg_{i}", use_container_width=True):
                        st.session_state.quick_q = q
                        st.rerun()
            scols2 = st.columns(3)
            for i, (col, q) in enumerate(zip(scols2, row2)):
                with col:
                    if st.button(q, key=f"sugg_{i+3}", use_container_width=True):
                        st.session_state.quick_q = q
                        st.rerun()
        else:
            # Render history
            for msg in st.session_state.chat_history:
                ts = msg.get("ts", "")
                content = msg["content"]

                if msg["role"] == "user":
                    safe = content.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br>")
                    st.markdown(
                        f'<div class="msg-group">'
                        f'<div class="msg-user-row">'
                        f'<div class="msg-user-bubble">{safe}</div>'
                        f'</div>'
                        f'<div class="msg-ts">{ts}</div>'
                        f'</div>',
                        unsafe_allow_html=True,
                    )
                else:
                    # Render assistant message via st.markdown for proper formatting
                    st.markdown(
                        '<div class="msg-group"><div class="msg-ai-row">'
                        '<div class="msg-ai-avatar">AI</div>'
                        '<div style="flex:1;max-width:72%">',
                        unsafe_allow_html=True,
                    )
                    st.markdown(
                        f'<div class="msg-ai-bubble">{content}</div>',
                        unsafe_allow_html=True,
                    )
                    # Sources (only if toggle is ON)
                    if st.session_state.get("show_sources") and msg.get("sources"):
                        tags = " ".join(
                            f'<span class="source-tag">{s[:30]}</span>'
                            for s in msg["sources"][:5]
                        )
                        st.markdown(
                            f'<div class="sources-bar">{tags}</div>',
                            unsafe_allow_html=True,
                        )
                    st.markdown(
                        f'<div class="msg-ts msg-ts-left">{ts}</div>'
                        '</div></div></div>',
                        unsafe_allow_html=True,
                    )

        # Scroll anchor
        st.markdown("<div id='chat-end'></div>", unsafe_allow_html=True)
        # Streaming response renders here while generating
        stream_ph = st.empty()
        st.markdown("<div style='height:130px'></div>", unsafe_allow_html=True)

    # ── Fixed input bar
    st.markdown('<div class="input-bar-wrap"><div class="input-bar-inner">', unsafe_allow_html=True)
    input_l, input_r = st.columns([9, 1])
    with input_l:
        user_input = st.text_input(
            "query", placeholder="Ask a question…",
            label_visibility="collapsed", key="chat_input",
            value=st.session_state.quick_q,
        )
    with input_r:
        send = st.button("Send", type="primary", use_container_width=True)
    st.markdown('</div></div>', unsafe_allow_html=True)

    # Clear quick_q after it was used
    if st.session_state.quick_q:
        st.session_state.quick_q = ""

    # ── Handle message send
    active = user_input.strip()
    if (send or (active and st.session_state._last_input != active)) and active:
        st.session_state._last_input = active
        ts_now = datetime.now().strftime("%H:%M")
        st.session_state.chat_history.append(
            {"role": "user", "content": active, "ts": ts_now}
        )
        answer, sources = ask_claude_stream(active, stream_ph)
        st.session_state.chat_history.append(
            {"role": "assistant", "content": answer, "ts": ts_now, "sources": sources}
        )
        st.rerun()


# ══════════════════════════════════════════════════════════════════
#  SETTINGS PAGE
# ══════════════════════════════════════════════════════════════════
def ts_label(iso):
    try: return datetime.fromisoformat(iso).strftime("%d %b %H:%M")
    except: return ""

def render_settings():
    render_topnav(show_settings_btn=False, show_back_btn=True)

    docs, links, wa, pages = kb_stats()

    # Page header
    st.markdown("""
    <div class="settings-page">
    <div class="settings-inner">
    """, unsafe_allow_html=True)

    # Back link + title row
    st.markdown("""
    <div class="settings-title">Knowledge Base Settings</div>
    <div class="settings-sub">
        Manage the documents and data sources that power the AI assistant.
        None of these sources are visible to end users.
    </div>
    """, unsafe_allow_html=True)

    # ── Stats row ──────────────────────────────────────────────────
    s1, s2, s3, s4, s5 = st.columns(5)
    s1.metric("Documents", docs)
    s2.metric("Web pages", links)
    s3.metric("WA chats",  wa)
    s4.metric("Crawled",   pages)
    s5.metric("Total",     docs + links + wa + pages)

    st.markdown("---")

    # ── Tabs ────────────────────────────────────────────────────────
    t1, t2, t3, t4, t5 = st.tabs([
        "📄 Documents", "🌐 Web Links", "💬 WhatsApp", "🕷️ Crawl Site", "⚙️ Preferences"
    ])

    # ── Documents ──────────────────────────────────────────────────
    with t1:
        st.caption("Supported: PDF · DOCX · XLSX · PPTX · HTML · TXT · CSV · JSON · RTF · EPUB · EML · Jupyter (.ipynb) · code files · and more.")
        ups = st.file_uploader(
            "upload_docs", accept_multiple_files=True,
            key="doc_uploader", label_visibility="collapsed",
        )
        if ups:
            ex = {d["name"] for d in st.session_state.kb_documents}
            added = 0
            for f in ups:
                if f.name not in ex:
                    with st.spinner(f"Parsing {f.name}…"):
                        content = parse_file(f)
                    st.session_state.kb_documents.append({
                        "name": f.name,
                        "content": content,
                        "type": _ext(f.name),
                        "added_at": datetime.now().isoformat(),
                        "size": len(content),
                    })
                    added += 1
            if added:
                save_kb()
                st.success(f"✅ {added} file(s) added and saved.")

        if st.session_state.kb_documents:
            st.markdown("**Loaded documents**")
            for i, d in enumerate(st.session_state.kb_documents):
                ca, cb = st.columns([6, 1])
                with ca:
                    st.markdown(
                        f'<div class="file-row">'
                        f'<span class="file-row-icon">{FILE_ICON(d["type"])}</span>'
                        f'<span class="file-row-name">{d["name"]}</span>'
                        f'<span class="file-row-meta">{d["size"]:,} chars · {ts_label(d["added_at"])}</span>'
                        f'</div>',
                        unsafe_allow_html=True,
                    )
                with cb:
                    if st.button("Remove", key=f"rm_doc_{i}", type="secondary"):
                        st.session_state.kb_documents.pop(i)
                        save_kb(); st.rerun()

    # ── Web Links ──────────────────────────────────────────────────
    with t2:
        st.caption("Paste a URL to scrape and index that page's content.")
        url_in = st.text_input("URL", placeholder="https://convin.ai/features", key="link_url")
        if st.button("Add page", key="add_link", type="primary"):
            u = url_in.strip()
            if u:
                if u in {l["url"] for l in st.session_state.kb_links}:
                    st.warning("This URL is already in the knowledge base.")
                else:
                    with st.spinner("Fetching page…"):
                        title, content = fetch_url(u)
                    st.session_state.kb_links.append({
                        "url": u, "title": title, "content": content,
                        "added_at": datetime.now().isoformat(), "size": len(content),
                    })
                    save_kb()
                    st.success(f"Added: {title[:50]}")
            else:
                st.error("Enter a URL first.")

        if st.session_state.kb_links:
            st.markdown("**Loaded links**")
            for i, l in enumerate(st.session_state.kb_links):
                ca, cb = st.columns([6, 1])
                with ca:
                    st.markdown(
                        f'<div class="file-row">'
                        f'<span class="file-row-icon">🌐</span>'
                        f'<span class="file-row-name">{l["title"][:50]}</span>'
                        f'<span class="file-row-meta">{l["size"]:,} chars · {ts_label(l["added_at"])}</span>'
                        f'</div>',
                        unsafe_allow_html=True,
                    )
                with cb:
                    if st.button("Remove", key=f"rm_link_{i}", type="secondary"):
                        st.session_state.kb_links.pop(i)
                        save_kb(); st.rerun()

    # ── WhatsApp ───────────────────────────────────────────────────
    with t3:
        st.markdown("""
        <div style='background:rgba(16,185,129,0.07);border:1px solid rgba(16,185,129,0.20);
        border-radius:12px;padding:12px 16px;margin-bottom:14px;font-size:0.82rem;color:#6EE7B7'>
        📱 <b>How to export:</b> Open WhatsApp → Chat → ⋮ / ⋯ → More → Export Chat → share the <code>.txt</code> or <code>.zip</code> file here.
        </div>
        """, unsafe_allow_html=True)

        wa_up = st.file_uploader(
            "wa_upload", accept_multiple_files=True,
            type=["txt", "zip", "pdf", "docx", "doc", "xlsx", "xls", "csv", "pptx", "ppt", "png", "jpg", "jpeg", "webp", "gif", "mp4", "mp3", "ogg", "opus", "aac", "wav", "json", "html", "htm"],
            key="wa_uploader", label_visibility="collapsed",
        )
        if wa_up:
            ex = {w["name"] for w in st.session_state.kb_whatsapp}
            added, skipped = 0, []
            for f in wa_up:
                if f.name in ex:
                    continue

                # Extract text from zip (WhatsApp export with media)
                if f.name.lower().endswith(".zip"):
                    try:
                        with zipfile.ZipFile(io.BytesIO(f.read())) as zf:
                            txt_files = [n for n in zf.namelist() if n.endswith(".txt")]
                            if not txt_files:
                                skipped.append(f.name)
                                continue
                            raw = zf.read(txt_files[0]).decode("utf-8", "ignore")
                            display_name = f.name
                    except Exception:
                        skipped.append(f.name)
                        continue
                elif f.name.lower().endswith(".txt"):
                    raw = f.read().decode("utf-8", "ignore")
                    display_name = f.name
                else:
                    # Non-text files: store as attachment reference
                    b64 = base64.b64encode(f.read()).decode()
                    st.session_state.kb_whatsapp.append({
                        "name": f.name, "content": f"[Attached file: {f.name}]",
                        "added_at": datetime.now().isoformat(),
                        "size": f.size,
                        "meta": {"valid": True, "total": 0, "participants": [], "date_range": "", "is_attachment": True},
                        "attachment_data": b64,
                        "attachment_mime": f.type or "application/octet-stream",
                    })
                    added += 1
                    continue

                parsed = parse_wa(raw)
                meta   = parse_wa_meta(parsed)

                # Validate it's a real WhatsApp export
                if not meta["valid"] or meta["total"] < 3:
                    skipped.append(f.name)
                    continue

                st.session_state.kb_whatsapp.append({
                    "name": display_name, "content": parsed,
                    "added_at": datetime.now().isoformat(),
                    "size": len(parsed),
                    "meta": meta,
                })
                added += 1

            if added:
                save_kb()
                st.success(f"✅ {added} file(s) added.")
            if skipped:
                st.warning(f"⚠️ {', '.join(skipped)} — doesn't look like a WhatsApp export (no messages found).")

        if st.session_state.kb_whatsapp:
            st.markdown("**Loaded chats**")
            for i, w in enumerate(st.session_state.kb_whatsapp):
                meta = w.get("meta") or parse_wa_meta(w.get("content", ""))
                plist = ", ".join(meta.get("participants", [])[:4]) if meta.get("valid") else "—"
                if len(meta.get("participants", [])) > 4:
                    plist += f" +{len(meta['participants'])-4} more"
                drange = meta.get("date_range", "")
                total_m = meta.get("total", 0)

                ca, cb = st.columns([6, 1])
                with ca:
                    st.markdown(
                        f'<div class="file-row" style="flex-direction:column;align-items:flex-start;gap:6px;padding:12px 16px">'
                        f'<div style="display:flex;align-items:center;gap:8px;width:100%">'
                        f'<span class="file-row-icon">💬</span>'
                        f'<span class="file-row-name">{w["name"]}</span>'
                        f'<span class="file-row-meta">{w["size"]:,} chars · {ts_label(w["added_at"])}</span>'
                        f'</div>'
                        f'<div style="display:flex;gap:10px;flex-wrap:wrap;padding-left:24px">'
                        f'<span style="font-size:0.7rem;color:#10B981">👥 {plist}</span>'
                        f'<span style="font-size:0.7rem;color:#8D90AA">📅 {drange}</span>'
                        f'<span style="font-size:0.7rem;color:#8D90AA">💬 {total_m:,} messages</span>'
                        f'</div>'
                        f'</div>',
                        unsafe_allow_html=True,
                    )
                with cb:
                    if st.button("Remove", key=f"rm_wa_{i}", type="secondary"):
                        st.session_state.kb_whatsapp.pop(i)
                        save_kb(); st.rerun()

                content = w.get("content", "")
                if content and not w.get("meta", {}).get("is_attachment"):
                    with st.expander(f"📖 View messages — {w['name']}", expanded=False):
                        msg_lines = [ln for ln in content.split("\n") if ln.strip()]
                        show_last = 300
                        if len(msg_lines) > show_last:
                            st.caption(f"Showing last {show_last} of {len(msg_lines):,} messages")
                            msg_lines = msg_lines[-show_last:]
                        st.markdown(
                            '<div style="height:420px;overflow-y:auto;font-family:monospace;'
                            'font-size:0.73rem;line-height:1.6;background:#0d1117;color:#c9d1d9;'
                            'padding:12px 16px;border-radius:8px;white-space:pre-wrap;">'
                            + "\n".join(msg_lines).replace("<", "&lt;").replace(">", "&gt;")
                            + "</div>",
                            unsafe_allow_html=True,
                        )

    # ── Crawl Site ─────────────────────────────────────────────────
    with t4:
        st.caption("Automatically crawl all pages on a website and add them to the knowledge base.")
        crawl_url = st.text_input("Root URL", value="https://convin.ai", key="crawl_url")
        max_p     = st.slider("Max pages", 5, 100, 30, 5)
        col_go, col_clr = st.columns([2, 1])
        go  = col_go.button("Start crawl", type="primary", use_container_width=True)
        clr = col_clr.button("Clear all crawled", type="secondary", use_container_width=True)

        if clr:
            st.session_state.kb_crawled = []; save_kb(); st.rerun()

        s_ph = st.empty(); p_ph = st.empty()
        if go:
            u = crawl_url.strip()
            if u:
                s_ph.info(f"Crawling **{u}** …")
                n = crawl_site(u, max_p, s_ph, p_ph)
                save_kb()
                s_ph.success(f"✅ Done — {n} pages indexed.")
                st.rerun()
            else:
                st.error("Enter a URL.")

        if st.session_state.kb_crawled:
            with st.expander(f"📑 {len(st.session_state.kb_crawled)} crawled pages", expanded=False):
                for i, pg in enumerate(st.session_state.kb_crawled):
                    ca, cb = st.columns([6, 1])
                    with ca:
                        st.markdown(
                            f'<div class="file-row">'
                            f'<span class="file-row-icon">🕷️</span>'
                            f'<span class="file-row-name">{pg["title"][:55]}</span>'
                            f'<span class="file-row-meta">{pg["size"]:,} chars</span>'
                            f'</div>',
                            unsafe_allow_html=True,
                        )
                    with cb:
                        if st.button("Remove", key=f"rm_pg_{i}", type="secondary"):
                            st.session_state.kb_crawled.pop(i)
                            save_kb(); st.rerun()

    # ── Preferences ────────────────────────────────────────────────
    with t5:
        st.markdown("**Display settings**")
        show_src = st.toggle(
            "Show sources under AI answers",
            value=st.session_state.get("show_sources", False),
            help="When enabled, the document/page names used to generate each answer are shown below the response.",
        )
        if show_src != st.session_state.get("show_sources", False):
            st.session_state["show_sources"] = show_src
            save_kb()

        st.markdown("---")
        st.markdown("**Export knowledge base**")
        _exp_faqs = st.session_state.get("kb_faqs", [])
        if _exp_faqs:
            _exp_cats = sorted(set(f["category"] for f in _exp_faqs))
            ex1, ex2 = st.columns(2)
            with ex1:
                st.download_button(
                    "⬇️ Export JSON",
                    data=json.dumps(_exp_faqs, indent=2, ensure_ascii=False),
                    file_name="answer-studio.json", mime="application/json",
                    use_container_width=True,
                )
            with ex2:
                _lines = []
                for _cat in _exp_cats:
                    _lines += [f"\n{'='*50}", f"  {_cat.upper()}", f"{'='*50}\n"]
                    for _idx, _faq in enumerate([f for f in _exp_faqs if f["category"] == _cat], 1):
                        _lines += [f"Q{_idx}. {_faq['question']}", f"A:  {_faq['answer']}\n"]
                st.download_button(
                    "⬇️ Export TXT",
                    data="\n".join(_lines), file_name="answer-studio.txt", mime="text/plain",
                    use_container_width=True,
                )
        else:
            st.caption("Generate answers first to enable exports.")

        st.markdown("---")
        st.markdown("**Danger zone**")
        if st.button("🗑️  Clear entire knowledge base", type="secondary"):
            for k in KB_KEYS:
                st.session_state[k] = []
            save_kb()
            st.success("Knowledge base cleared.")
            st.rerun()

    st.markdown("</div></div>", unsafe_allow_html=True)


# ══════════════════════════════════════════════════════════════════
#  FAQ PAGE  —  3-tab layout
#  Tab 1: FAQ       — curated Convin Sense product Q&As only
#  Tab 2: Q&A       — all generic Q&As (docs + web, no WhatsApp)
#  Tab 3: WhatsApp  — WhatsApp-extracted Q&As only
# ══════════════════════════════════════════════════════════════════

# Categories that are specifically about the Convin Sense product.
# Only these appear in the curated FAQ tab.
_CONVIN_FAQ_CATS = {
    # Core platform
    "Convin Platform",
    "AI Phone Calls",
    "AI Phone Call Agent",
    "Real-Time Assist Agent",
    "Conversation Intelligence",
    "Auto QA",
    "Convin Insights",
    "Convin LMS",
    # Product descriptions
    "Core Capabilities",
    "Core Features",
    "Product Overview",
    "Platform Overview",
    "How It Works",
    "Key Value Propositions",
    "Getting Started",
    # Business & trust
    "Integrations",
    "Data Security",
    "Human-in-the-Loop",
    "Business Value",
    "Target Audience",
    "Platform Positioning",
    "Use Cases",
    "Trust & Reputation",
    "Call Quality & QA",
    "Definitions & Concepts",
    "Company & Background",
    "Pricing & Plans",
    # Use cases
    "Sales Use Case",
    "Collections Use Case",
    "Compliance Use Case",
    "Lead Qualification",
    "Customer Retention",
}
# Max Q&As shown per category in the FAQ tab to keep it concise
_FAQ_CAP_PER_CAT = 5

_CAT_ICONS = {
    # WhatsApp pilot categories
    "WhatsApp: Pilot Clients & Metrics":        "📊",
    "WhatsApp: Product Issues & Bugs":          "🐛",
    "WhatsApp: Setup & Configuration":          "⚙️",
    "WhatsApp: Onboarding Learnings":           "🚀",
    "WhatsApp: Client Objections & Responses":  "💬",
    "WhatsApp: Feature Requests":               "✨",
    "WhatsApp: Bot Performance":                "🤖",
    "WhatsApp: Sales Process":                  "💼",
    # Product / FAQ categories
    "AI Phone Call Platform":           "📞",
    "Conversation Intelligence":        "🧠",
    "Real-Time Assist":                 "⚡",
    "Auto QA":                          "✅",
    "Convin Insights":                  "📈",
    "Pricing & Plans":                  "💰",
    "Integrations":                     "🔗",
    "Company & Background":             "🏢",
    "Sales Use Case":                   "💼",
    "Coaching":                         "🎯",
    "LMS":                              "📚",
    "Healthcare":                       "🏥",
    "Banking & Finance":                "🏦",
    "Compliance":                       "🛡️",
    "Customer Retention":               "🤝",
    "Lead Qualification":               "🎯",
    "Collections":                      "💳",
    "BPO":                              "🏭",
    "Insurance":                        "📋",
    "Home Services":                    "🏠",
    # Generic
    "Business Outcomes & ROI":          "📊",
    "Core Capabilities":                "⚙️",
    "How It Works":                     "🔍",
    "After Hours Answering Service":    "🌙",
    "Answering Service Pricing Models": "💰",
    "AI Call Answering Service":        "🤖",
    "Customer Satisfaction Score":      "⭐",
    "Call Queue Management":            "📋",
    "Bot Performance":                  "🤖",
    "Feature Requests":                 "✨",
    "Setup & Configuration":            "⚙️",
    "Onboarding Learnings":             "🚀",
    "Sales Process":                    "💼",
}

def _render_category_dashboard(subset: list[dict], tab_key: str, no_content_msg: str = "No Q&As yet."):
    """Category cards + per-category expandable Q&A list."""
    if not subset:
        st.markdown(f"""
        <div class="no-faq">
          <div class="no-faq-icon">✦</div>
          <h3>{no_content_msg}</h3>
          <p>Generate answers first using the button above.</p>
        </div>
        """, unsafe_allow_html=True)
        return

    cat_counts = {}
    for f in subset:
        cat_counts[f["category"]] = cat_counts.get(f["category"], 0) + 1

    # ── Stat cards row ────────────────────────────────────────────
    cards_html = ""
    for cat, count in sorted(cat_counts.items(), key=lambda x: -x[1]):
        icon = _CAT_ICONS.get(cat, "📌")
        label = cat.replace("WhatsApp: ", "")
        cards_html += (
            f'<div style="background:#1a1d2e;border:1px solid #2d3158;border-radius:10px;'
            f'padding:14px 16px;min-width:150px;flex:1;max-width:220px">'
            f'<div style="font-size:1.3rem">{icon}</div>'
            f'<div style="font-size:1.5rem;font-weight:700;color:#A78BFA;margin:4px 0">{count}</div>'
            f'<div style="font-size:0.72rem;color:#9CA3AF;line-height:1.4">{label}</div>'
            f'</div>'
        )
    st.markdown(
        f'<div style="display:flex;flex-wrap:wrap;gap:10px;margin:16px 0 24px">{cards_html}</div>',
        unsafe_allow_html=True,
    )

    # ── Per-category expanders ────────────────────────────────────
    for cat, count in sorted(cat_counts.items(), key=lambda x: -x[1]):
        icon = _CAT_ICONS.get(cat, "📌")
        label = cat.replace("WhatsApp: ", "")
        bucket = [f for f in subset if f["category"] == cat]
        safe_key = cat.replace(" ", "_").replace(":", "")[:25]
        with st.expander(f"{icon}  {label}  ·  {count} Q&As", expanded=False):
            _render_faq_list(bucket, f"{tab_key}_{safe_key}", f"search_{tab_key}_{safe_key}")

def _render_faq_list(subset: list[dict], tab_key: str, search_key: str):
    """Reusable search + category expanders for a subset of FAQs."""
    if not subset:
        st.markdown("""
        <div class="no-faq">
          <div class="no-faq-icon">✦</div>
          <h3>Nothing here yet</h3>
          <p>Generate answers first using the button above.</p>
        </div>
        """, unsafe_allow_html=True)
        return

    cats = list(dict.fromkeys(f["category"] for f in subset))

    sc, qc = st.columns([6, 1])
    with sc:
        search = st.text_input(
            "search", placeholder="🔍  Search questions and answers…",
            label_visibility="collapsed", key=search_key,
        )
    with qc:
        st.markdown(
            f"<div style='text-align:right;padding-top:8px;"
            f"font-size:0.78rem;color:#8D90AA'>"
            f"<b style='color:#A78BFA'>{len(subset)}</b> Q&As</div>",
            unsafe_allow_html=True,
        )
    slc = search.lower().strip() if search else ""

    for cat in cats:
        cat_faqs = [f for f in subset if f["category"] == cat]
        if slc:
            cat_faqs = [
                f for f in cat_faqs
                if slc in f["question"].lower() or slc in f["answer"].lower()
            ]
        if not cat_faqs:
            continue

        st.markdown(
            f'<div class="cat-label">📂 {cat} &nbsp;·&nbsp; {len(cat_faqs)} questions</div>',
            unsafe_allow_html=True,
        )

        for idx, faq in enumerate(cat_faqs):
            q_disp = faq["question"]
            a_disp = faq["answer"]

            has_wa = "💬 Chatted by" in a_disp or "chatted by" in a_disp.lower()
            badge_html = (
                '<span class="faq-wa-badge">💬 WhatsApp source</span>'
                if has_wa else
                '<span class="faq-doc-badge">📄 KB source</span>'
            )

            if slc:
                def hl(text, term=slc):
                    return re.sub(
                        f"({re.escape(term)})",
                        r'<mark style="background:rgba(124,58,237,0.28);'
                        r'border-radius:3px;padding:0 3px;color:#EEF0FA">\1</mark>',
                        text, flags=re.IGNORECASE,
                    )
                a_disp = hl(a_disp)
                q_disp = hl(q_disp)

            a_rendered = re.sub(
                r"(💬 Chatted by[^\n]+)",
                r'<div class="wa-cite">🟢 \1</div>',
                a_disp,
            )

            with st.expander(f"Q: {faq['question']}", expanded=False):
                st.markdown(
                    f"<div class='faq-answer-wrap'>"
                    f"<div style='margin-bottom:10px'>{badge_html}</div>"
                    f"<div class='faq-answer-label'>Answer</div>"
                    f"<div class='faq-answer-body'>{a_rendered}</div>"
                    f"</div>",
                    unsafe_allow_html=True,
                )
                if st.button(
                    "💬 Ask follow-up in chat",
                    key=f"faq_ask_{tab_key}_{cat}_{idx}",
                    type="secondary",
                ):
                    st.session_state.quick_q  = faq["question"]
                    st.session_state.chat_open = True
                    st.rerun()


def _render_mini_chat():
    """Inline chat panel rendered as a side column on the FAQ page."""
    # ── Panel header
    st.markdown("""
    <div class="chat-panel-wrap">
    <div class="chat-panel-header">
      <div class="chat-panel-title">
        <span class="live-dot-sm"></span>
        <div>
          AI Chat
          <div class="sub">Animesh · Powered by Claude</div>
        </div>
      </div>
    </div>
    """, unsafe_allow_html=True)

    # ── Close button
    if st.button("✕ Close", key="close_chat_panel", type="secondary"):
        st.session_state.chat_open = False
        st.session_state.quick_q   = ""
        st.rerun()

    # ── Messages feed
    history = st.session_state.get("chat_history", [])
    msgs_html = ""
    for msg in history[-20:]:
        content = msg["content"]
        if msg["role"] == "user":
            safe = content.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br>")
            msgs_html += f'<div class="mini-msg-user"><span>{safe}</span></div>'
        else:
            msgs_html += (
                f'<div class="mini-msg-ai">'
                f'<div class="av">AI</div>'
                f'<div class="bubble">{content}</div>'
                f'</div>'
            )

    if msgs_html:
        st.markdown(
            f'<div class="chat-panel-msgs">{msgs_html}</div>',
            unsafe_allow_html=True,
        )
    else:
        st.markdown("""
        <div class="chat-panel-msgs">
          <div class="mini-chat-empty">
            <div class="icon">✦</div>
            <div class="text">
              Hi, I'm Animesh.<br>Ask me anything about Convin Sense.
            </div>
          </div>
        </div>
        """, unsafe_allow_html=True)

    st.markdown("</div>", unsafe_allow_html=True)

    # ── Streaming placeholder (visible while generating)
    stream_ph = st.empty()

    # ── Input — st.chat_input handles Enter natively, no button needed
    user_input = st.chat_input("Ask a question…", key="mini_chat_input")

    # Handle quick_q (suggestion chips)
    if not user_input and st.session_state.get("quick_q"):
        user_input = st.session_state.quick_q
        st.session_state.quick_q = ""

    if user_input and user_input.strip():
        active = user_input.strip()
        ts_now = datetime.now().strftime("%H:%M")
        st.session_state.chat_history.append({"role": "user", "content": active, "ts": ts_now})
        answer, sources = ask_claude_stream(active, stream_ph)
        st.session_state.chat_history.append(
            {"role": "assistant", "content": answer, "ts": ts_now, "sources": sources}
        )
        st.session_state["_mini_last"] = active
        st.rerun()

    if st.button("🗑 Clear chat", key="mini_clear_chat", type="secondary", use_container_width=True):
        st.session_state.chat_history = []
        st.session_state["_mini_last"] = ""
        st.rerun()


def _float_via_js(marker_id: str, width: str, bottom: str = "24px", right: str = "24px", z: str = "9999"):
    """Inject JS that positions the Streamlit container holding `marker_id` as fixed overlay."""
    st.markdown(f"""<script>
(function(){{
    var F=function(){{
        var m=document.getElementById('{marker_id}');
        if(!m)return;
        var el=m;
        for(var i=0;i<14;i++){{
            el=el.parentElement;
            if(!el)break;
            var t=el.getAttribute&&el.getAttribute('data-testid');
            if(t==='stVerticalBlock'){{
                el.style.setProperty('position','fixed','important');
                el.style.setProperty('bottom','{bottom}','important');
                el.style.setProperty('right','{right}','important');
                el.style.setProperty('width','{width}','important');
                el.style.setProperty('z-index','{z}','important');
                el.style.setProperty('margin','0','important');
                el.style.setProperty('padding','0','important');
                break;
            }}
        }}
    }};
    F();setTimeout(F,80);setTimeout(F,300);
}})();
</script>""", unsafe_allow_html=True)


def _render_chat_float():
    """Floating chat overlay — stays on the landing page, no layout shift."""
    chat_open = st.session_state.get("chat_open", False)

    if not chat_open:
        # ── FAB only ──────────────────────────────────────────────
        st.markdown('<div id="cf-fab-mk" style="display:none"></div>', unsafe_allow_html=True)
        st.markdown('<div class="cf-fab-wrap">', unsafe_allow_html=True)
        if st.button("💬", key="cf_fab_open", help="Chat with Animesh"):
            st.session_state.chat_open = True
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)
        _float_via_js("cf-fab-mk", "64px")
        return

    # ── Full chat panel ───────────────────────────────────────────
    st.markdown('<div id="cf-panel-mk" style="display:none"></div>', unsafe_allow_html=True)
    st.markdown('<div class="cf-shell">', unsafe_allow_html=True)

    # Header
    st.markdown(
        '<div class="cf-header">'
        '<div class="cf-title-row">'
        '<span class="cf-live-dot"></span>'
        '<div><div class="cf-title">AI Chat</div>'
        '<div class="cf-sub">Animesh · Powered by Claude</div></div>'
        '</div>'
        '</div>',
        unsafe_allow_html=True,
    )

    # Close button (right-aligned)
    st.markdown('<div style="text-align:right;padding:6px 12px 0" class="cf-close-btn">', unsafe_allow_html=True)
    if st.button("✕ Close", key="cf_close", type="secondary"):
        st.session_state.chat_open = False
        st.session_state.quick_q = ""
        st.rerun()
    st.markdown('</div>', unsafe_allow_html=True)

    # Messages
    history = st.session_state.get("chat_history", [])
    if history:
        msgs = ""
        for msg in history[-20:]:
            c = msg["content"]
            if msg["role"] == "user":
                safe = c.replace("&","&amp;").replace("<","&lt;").replace(">","&gt;").replace("\n","<br>")
                msgs += f'<div class="cf-user-bubble">{safe}</div>'
            else:
                msgs += f'<div class="cf-ai-bubble">{c}</div>'
        st.markdown(f'<div class="cf-msgs">{msgs}</div>', unsafe_allow_html=True)
    else:
        st.markdown(
            '<div class="cf-msgs"><div class="cf-empty">'
            '<span class="cf-empty-icon">✦</span>'
            "Hi, I'm Animesh.<br>Ask me anything about Convin Sense."
            '</div></div>',
            unsafe_allow_html=True,
        )

    st.markdown('</div>', unsafe_allow_html=True)  # close cf-shell

    # ── Streaming placeholder ─────────────────────────────────────
    stream_ph = st.empty()

    # ── Input row ────────────────────────────────────────────────
    st.markdown('<div class="cf-input-row">', unsafe_allow_html=True)
    in_col, send_col = st.columns([5, 1])
    with in_col:
        pre = st.session_state.get("quick_q", "")
        user_input = st.text_input(
            "cf_msg", placeholder="Ask a question…",
            label_visibility="collapsed", key="cf_text_input",
            value=pre,
        )
        if pre:
            st.session_state.quick_q = ""
    with send_col:
        send_btn = st.button("↵", key="cf_send", type="primary")
    st.markdown('</div>', unsafe_allow_html=True)

    # Handle send
    active = user_input.strip() if user_input else ""
    if (send_btn or active) and active and active != st.session_state.get("_mini_last", ""):
        st.session_state["_mini_last"] = active
        ts = datetime.now().strftime("%H:%M")
        st.session_state.chat_history.append({"role": "user", "content": active, "ts": ts})
        answer, sources = ask_claude_stream(active, stream_ph)
        st.session_state.chat_history.append({"role": "assistant", "content": answer, "ts": ts, "sources": sources})
        st.rerun()

    # Clear button
    if history:
        st.markdown('<div style="padding:0 12px 10px">', unsafe_allow_html=True)
        if st.button("🗑 Clear chat", key="cf_clear", type="secondary", use_container_width=True):
            st.session_state.chat_history = []
            st.session_state["_mini_last"] = ""
            st.rerun()
        st.markdown('</div>', unsafe_allow_html=True)

    # JS: position this block fixed in bottom-right
    _float_via_js("cf-panel-mk", "370px", bottom="90px")


def render_faq():
    render_topnav(show_settings_btn=False, show_back_btn=False, show_chat_btn=True)

    faqs  = st.session_state.get("kb_faqs", [])
    total = total_sources()

    # Partition into WhatsApp, generic, and curated Convin FAQ
    wa_faqs      = [f for f in faqs if f["category"].startswith("WhatsApp:")]
    generic_faqs = [f for f in faqs if not f["category"].startswith("WhatsApp:")]

    # Curated FAQ: only product-specific categories, capped per category
    faq_curated: list[dict] = []
    for cat in _CONVIN_FAQ_CATS:
        bucket = [f for f in faqs if f["category"] == cat]
        faq_curated.extend(bucket[:_FAQ_CAP_PER_CAT])

    all_cats = list(dict.fromkeys(f["category"] for f in faqs)) if faqs else []
    wa_cats  = list(dict.fromkeys(f["category"] for f in wa_faqs))

    # ── Main content — always full-width ──────────────────────────
    with st.container():

    # ── Hero ──────────────────────────────────────────────────────
        # ── Landing Hero ──────────────────────────────────────────────
        st.markdown(
            f'<div class="lp-hero"><div class="lp-hero-inner">'
            f'<div class="lp-eyebrow"><span class="dot"></span>Convin Sense &nbsp;&middot;&nbsp; AI Knowledge Platform</div>'
            f'<div class="lp-title">Your Intelligent <span class="grad">Answer Studio</span></div>'
            f'<div class="lp-sub">Auto-generated answers from your entire knowledge base — docs, web pages, and conversations — always up to date.</div>'
            f'<div class="lp-pills">'
            f'<span class="lp-pill lp-pill-violet">&#10022; {len(faqs):,} Q&amp;As</span>'
            f'<span class="lp-pill lp-pill-pink">&#129302; AI-Powered</span>'
            f'<span class="lp-pill lp-pill-cyan">&#9889; Instant Search</span>'
            f'<span class="lp-pill lp-pill-green">&#10003; {total} KB Sources</span>'
            f'<span class="lp-pill lp-pill-amber">&#128172; Chat Ready</span>'
            f'</div>'
            f'<div class="lp-stats">'
            f'<div class="lp-stat lp-stat-v"><span class="n">{len(faq_curated)}</span><div class="l">Curated FAQ</div></div>'
            f'<div class="lp-stat lp-stat-p"><span class="n">{len(generic_faqs)}</span><div class="l">Q&amp;A Pairs</div></div>'
            f'<div class="lp-stat lp-stat-c"><span class="n">{len(wa_faqs)}</span><div class="l">WhatsApp Q&amp;A</div></div>'
            f'<div class="lp-stat lp-stat-g"><span class="n">{total}</span><div class="l">KB Sources</div></div>'
            f'</div>'
            f'</div></div>',
            unsafe_allow_html=True,
        )

        # ── Action bar ────────────────────────────────────────────────
        ab1, ab2 = st.columns([3, 2])
        with ab1:
            gen_btn = st.button(
                "✨ Generate Answers" if not faqs else "🔄 Regenerate Answers",
                type="primary", use_container_width=True, disabled=(total == 0),
            )
        with ab2:
            if faqs and st.button("🗑️ Clear All", use_container_width=True):
                st.session_state.kb_faqs = []
                save_kb(); st.rerun()
        if total == 0:
            st.info("Add documents or links from Settings first, then click Generate.")

        # ── Generate ──────────────────────────────────────────────────
        if gen_btn:
            prog_ph   = st.empty()
            status_ph = st.empty()

            docs_n  = len(st.session_state.get("kb_documents", []))
            links_n = len(st.session_state.get("kb_links", [])) + len(st.session_state.get("kb_crawled", []))
            wa_n    = len(st.session_state.get("kb_whatsapp", []))
            passes  = sum(1 for x in [docs_n, links_n, wa_n] if x > 0)

            status_ph.markdown(
                f"<span style='color:#A78BFA;font-size:0.85rem'>"
                f"🤖 Running {passes} extraction pass(es) across all sources — "
                f"building your Answer Studio…</span>",
                unsafe_allow_html=True,
            )
            prog_ph.progress(0.05)

            def _progress(pct, label):
                prog_ph.progress(pct)
                status_ph.markdown(
                    f"<span style='color:#A78BFA;font-size:0.85rem'>{label}</span>",
                    unsafe_allow_html=True,
                )

            try:
                new_faqs = generate_faqs(progress_cb=_progress)
                prog_ph.progress(1.0)
                if new_faqs:
                    st.session_state.kb_faqs = new_faqs
                    save_kb()
                    prog_ph.empty()
                    n_cats = len(set(f["category"] for f in new_faqs))
                    status_ph.success(
                        f"✅ {len(new_faqs)} answers generated across {n_cats} categories — saved to Answer Studio!"
                    )
                    st.rerun()
                else:
                    prog_ph.empty()
                    status_ph.warning("No answers extracted. Try adding more content.")
            except Exception as e:
                prog_ph.empty()
                status_ph.error(f"Error: {e}")

        # ── 3-Tab layout ──────────────────────────────────────────────
        tab_faq, tab_generic, tab_wa = st.tabs([
            f"📋  FAQ  ({len(faq_curated)})",
            f"💡  Q&A  ({len(generic_faqs)})",
            f"💬  WhatsApp  ({len(wa_faqs)})",
        ])

        # ── Tab 1: Curated Convin Sense FAQ ───────────────────────────
        with tab_faq:
            _render_category_dashboard(
                faq_curated, "faq",
                no_content_msg="No Convin Sense FAQs found",
            )

        # ── Tab 2: Full generic Q&A (docs + web) ─────────────────────
        with tab_generic:
            _render_category_dashboard(
                generic_faqs, "generic",
                no_content_msg="No Q&As yet — add docs or web links in Settings",
            )

        # ── Tab 3: WhatsApp Q&As ──────────────────────────────────────
        with tab_wa:
            _render_category_dashboard(
                wa_faqs, "wa",
                no_content_msg="No WhatsApp Q&As yet — upload a chat export in Settings",
            )

    # ── Floating chat widget (always rendered, positioned via JS) ──
    _render_chat_float()


# ══════════════════════════════════════════════════════════════════
#  ROUTER
# ══════════════════════════════════════════════════════════════════
if st.session_state.page == "chat":
    render_chat()
elif st.session_state.page == "settings":
    render_settings()
elif st.session_state.page == "faq":
    render_faq()
